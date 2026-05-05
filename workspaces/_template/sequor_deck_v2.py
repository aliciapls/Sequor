"""Generate Sequor PE Investor Presentation Deck — Professional Version."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x4A, 0x63)   # D800 – primary dark
TEAL    = RGBColor(0x3C, 0x8E, 0xAF)   # D500 – accent
TEAL_L  = RGBColor(0x5F, 0xA6, 0xC5)   # D400 – light accent
TEAL_XL = RGBColor(0xB7, 0xDD, 0xEA)   # D200 – bg accent
BG      = RGBColor(0xF0, 0xF9, 0xFB)   # D50 – slide bg
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BODY    = RGBColor(0x47, 0x56, 0x6A)   # slate text
MUTED   = RGBColor(0x94, 0xA3, 0xB8)   # muted text
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
RED     = RGBColor(0xEF, 0x44, 0x44)
DARK    = RGBColor(0x0F, 0x17, 0x2A)   # near black for cover

# ── Slide setup ───────────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]


# ── Drawing helpers ────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    from pptx.util import Emu
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.line.width = lw
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s


def tx(slide, text, l, t, w=None, h=None,
       size=Pt(18), bold=False, color=BODY,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    if w is None: w = Inches(1)
    if h is None: h = Inches(0.5)
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Inter"
    return tb


def tx_multi(slide, lines, l, t, w, line_h=Pt(22),
             size=Pt(14), color=BODY, bold_first=False):
    """Multi-line text box. lines is list of (text, bold, color)."""
    from pptx.util import Pt
    total_h = (len(lines) - 1) * line_h + Pt(28)
    tb = slide.shapes.add_textbox(l, t, w, total_h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, bold, c) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.size = size
        r.font.bold = bold
        r.font.color.rgb = c
        r.font.name = "Inter"
    return tb


def add_line(slide, l, t, w, h=Inches(0.025), color=TEAL):
    rect(slide, l, t, w, h, fill=color)


def card(slide, l, t, w, h, fill=WHITE, border=TEAL_XL, border_w=Pt(1)):
    rect(slide, l, t, w, h, fill=fill, line=border, lw=border_w)


def header_bar(slide, title, subtitle=None, dark=False):
    bg = DARK if dark else NAVY
    rect(slide, 0, 0, W, Inches(1.4), fill=bg)
    tx(slide, title, Inches(0.65), Inches(0.32), size=Pt(32), bold=True,
       color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        tx(slide, subtitle, Inches(0.65), Inches(0.92), size=Pt(14),
           color=TEAL_L, italic=True)


def slide_bg(slide, color=BG):
    rect(slide, 0, 0, W, H, fill=color)


def footer(slide, text="Confidential  •  May 2026  •  Sequor"):
    rect(slide, 0, Inches(7.1), W, Inches(0.4), fill=NAVY)
    tx(slide, text, Inches(0.65), Inches(7.15), size=Pt(10),
       color=TEAL_L, italic=True)


def metric(slide, l, t, value, label, accent=TEAL):
    """Big metric number with label below."""
    tx(slide, value, l, t, Inches(1.8), Inches(0.9),
       size=Pt(52), bold=True, color=accent, align=PP_ALIGN.CENTER)
    tx(slide, label, l, t + Inches(0.85), Inches(1.8), Inches(0.4),
       size=Pt(11), color=BODY, align=PP_ALIGN.CENTER)


def icon_bullet(slide, l, t, w, icon, title, desc, accent=TEAL):
    """Bullet with colored left bar, title bold, desc muted."""
    rect(slide, l, t, Inches(0.06), Inches(0.55), fill=accent)
    tx(slide, title, l + Inches(0.18), t, Inches(1.5), Inches(0.3),
       size=Pt(13), bold=True, color=NAVY)
    tx(slide, desc, l + Inches(0.18), t + Inches(0.28), w - Inches(0.18),
       size=Pt(11), color=BODY, wrap=True)


def pill(slide, l, t, text, bg=TEAL, fg=WHITE):
    """Small colored label pill."""
    card(slide, l, t, Inches(1.3), Inches(0.35), fill=bg, border=None)
    tx(slide, text, l + Inches(0.1), t + Inches(0.06),
       Inches(1.1), Inches(0.28),
       size=Pt(10), bold=True, color=fg, align=PP_ALIGN.CENTER)


# ── SLIDE 1: Cover ─────────────────────────────────────────────────────────────
def s1_cover():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, fill=DARK)
    # Right decorative panel
    rect(s, Inches(8.5), 0, Inches(4.83), H, fill=NAVY)
    rect(s, Inches(9.8), 0, Inches(3.53), H, fill=RGBColor(0x1D, 0x5A, 0x77))

    # Logo
    card(s, Inches(0.8), Inches(1.5), Inches(0.85), Inches(0.85), fill=TEAL)
    tx(s, "S", Inches(0.8), Inches(1.52), Inches(0.85), Inches(0.85),
       size=Pt(38), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, "Sequor", Inches(1.85), Inches(1.55), Inches(4), Inches(0.8),
       size=Pt(38), bold=True, color=WHITE)

    # Tagline
    add_line(s, Inches(0.8), Inches(2.55), Inches(6.5), color=TEAL)
    tx(s, "AI-Powered Customer Communication Platform",
       Inches(0.8), Inches(2.68), Inches(7), Inches(0.5),
       size=Pt(16), color=TEAL_L)

    # Headline
    tx(s, "The operating system for\ncustomer conversations.",
       Inches(0.8), Inches(3.5), Inches(7), Inches(1.6),
       size=Pt(38), bold=True, color=WHITE)

    # Sub
    tx(s, "One platform. Every channel. Fully automated.\nBuilt for Southeast Asian SMBs — PDPA compliant from day one.",
       Inches(0.8), Inches(5.3), Inches(7), Inches(0.8),
       size=Pt(14), color=MUTED)

    # Right panel text
    tx(s, "Investor\nPresentation",
       Inches(9.0), Inches(2.5), Inches(3.5), Inches(1.5),
       size=Pt(26), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, "May 2026", Inches(9.0), Inches(4.1), Inches(3.5), Inches(0.4),
       size=Pt(14), color=TEAL_L, align=PP_ALIGN.CENTER)
    tx(s, "Sequor Pte. Ltd.\nSingapore",
       Inches(9.0), Inches(5.5), Inches(3.5), Inches(0.7),
       size=Pt(12), color=MUTED, align=PP_ALIGN.CENTER)

    footer(s)


# ── SLIDE 2: Problem ────────────────────────────────────────────────────────────
def s2_problem():
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    header_bar(s, "The Problem",
              "Why businesses are losing customers — and operators — every day")
    footer(s)

    problems = [
        (RED,   "Channels in silos",
         "WhatsApp, Email, SMS — each a separate inbox. Context lost at every handoff."),
        (AMBER, "Response times kill loyalty",
         "Average first response: 4–6 hours. 35% of customers leave after one slow reply."),
        (RED,   "Repetitive queries burn out teams",
         "70% of inbound messages are repeats — FAQs, order status, product info."),
        (AMBER, "Escalations slip through",
         "High-intent signals (frustration, VIPs) are missed. Resolving one takes 6+ tools."),
        (RED,   "Compliance is manual and risky",
         "PDPA requires audit trails, erasure, data minimization. Manual processes fail audits."),
        (MUTED, "Zero operational visibility",
         "No unified view of volume, resolution rates, or AI accuracy. Data lives in inboxes."),
    ]

    cols = 3
    cw = Inches(3.8)
    ch = Inches(1.6)
    sx = Inches(0.65)
    sy = Inches(1.7)
    gap_x = Inches(0.28)
    gap_y = Inches(0.22)

    for i, (col, title, desc) in enumerate(problems):
        col_idx = i % cols
        row = i // cols
        x = sx + col_idx * (cw + gap_x)
        y = sy + row * (ch + gap_y)
        card(s, x, y, cw, ch)
        rect(s, x, y, Inches(0.08), ch, fill=col)
        tx(s, title, x + Inches(0.22), y + Inches(0.2), cw - Inches(0.35),
           size=Pt(14), bold=True, color=NAVY)
        tx(s, desc, x + Inches(0.22), y + Inches(0.55), cw - Inches(0.35),
           size=Pt(12), color=BODY, wrap=True)


# ── SLIDE 3: Solution ──────────────────────────────────────────────────────────
def s3_solution():
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    header_bar(s, "The Solution",
              "One platform — every channel, every message, fully automated")
    footer(s)

    # Central statement
    card(s, Inches(0.65), Inches(1.75), Inches(7.8), Inches(1.7))
    tx(s,
       "Sequor unifies WhatsApp and Email into a single AI-powered platform — auto-replies to common queries, "
       "intelligently escalates what matters, and gives operators complete visibility — while remaining "
       "PDPA compliant by design.",
       Inches(0.9), Inches(1.95), Inches(7.3), Inches(1.4),
       size=Pt(14), color=BODY, wrap=True)

    # 3 pillars
    pillars = [
        (TEAL,   "Auto-Reply AI",
         "Classifies intent, retrieves relevant documents via RAG, and generates accurate replies — without human input."),
        (NAVY,   "Smart Escalations",
         "Detects high-intent signals and routes to the right operator with full conversation context."),
        (TEAL_L, "PDPA Compliance",
         "Built-in audit trail, data minimization, and erasure workflows — compliant from day one."),
    ]
    pw = Inches(3.8)
    px = Inches(0.65)
    py = Inches(3.7)
    for col, title, desc in pillars:
        card(s, px, py, pw, Inches(1.65))
        rect(s, px, py, pw, Inches(0.08), fill=col)
        tx(s, title, px + Inches(0.25), py + Inches(0.22), pw - Inches(0.4),
           size=Pt(16), bold=True, color=NAVY)
        tx(s, desc, px + Inches(0.25), py + Inches(0.65), pw - Inches(0.4),
           size=Pt(12), color=BODY, wrap=True)
        px += pw + Inches(0.28)

    # Metrics row
    metrics = [("80%", "messages auto-replied"), ("4-hr", "avg SLA achievement"),
               ("< 10 min", "to first AI reply"), ("0", "servers to manage")]
    mx = Inches(0.65)
    my = Inches(5.6)
    for val, lbl in metrics:
        tx(s, val, mx, my, Inches(2.8), Inches(0.6),
           size=Pt(36), bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        tx(s, lbl, mx, my + Inches(0.6), Inches(2.8), Inches(0.4),
           size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)
        mx += Inches(3.1)


# ── SLIDE 4: Product ──────────────────────────────────────────────────────────
def s4_product():
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    header_bar(s, "Product: The Operator Portal",
              "Your command center for every customer interaction — one view, all channels")
    footer(s)

    # Left: feature list
    features = [
        (TEAL,   "Dashboard",     "Live metrics: message volume, auto-reply rate, SLA health, escalation queue"),
        (NAVY,   "Messages",      "Unified thread history across WhatsApp and Email — full context at a glance"),
        (TEAL,   "Escalations",   "Intelligent queue ranked by priority — assign, resolve, track resolution time"),
        (NAVY,   "Auto-Replies",  "Log of every AI-generated reply with confidence score and operator feedback"),
        (TEAL,   "Document Hub",  "Knowledge base linked to AI RAG — upload docs, AI uses them in every reply"),
        (NAVY,   "Key Phrases",   "Custom trigger phrases mapped to AI reply strategies per topic or channel"),
        (TEAL,   "Channels",      "WhatsApp and Email connections — webhook status, health, configuration"),
        (NAVY,   "Subscription",  "Plan, usage, and billing — transparent, predictable, self-serve"),
    ]
    fx = Inches(0.65)
    fy = Inches(1.75)
    fw = Inches(5.2)
    for col, title, desc in features:
        rect(s, fx, fy + Inches(0.06), Inches(0.06), Inches(0.42), fill=col)
        tx(s, title, fx + Inches(0.2), fy, Inches(1.5), Inches(0.28),
           size=Pt(12), bold=True, color=NAVY)
        tx(s, desc, fx + Inches(0.2), fy + Inches(0.28), fw - Inches(0.2),
           size=Pt(10.5), color=BODY, wrap=True)
        rect(s, fx, fy + Inches(0.62), fw, Inches(0.012), fill=TEAL_XL)
        fy += Inches(0.66)

    # Right: portal mockup frame
    card(s, Inches(6.3), Inches(1.65), Inches(6.6), Inches(5.1),
         fill=WHITE, border=NAVY, border_w=Pt(1.5))
    # Portal header bar
    rect(s, Inches(6.3), Inches(1.65), Inches(6.6), Inches(0.5), fill=NAVY)
    tx(s, "  Sequor  Portal", Inches(6.3), Inches(1.7), Inches(4),
       size=Pt(11), bold=True, color=WHITE)

    # Stat cards in mock
    stats = [("Messages", "1,284", "+12%", TEAL), ("Auto-Replied", "1,027", "80%", GREEN),
             ("Escalated", "42", "-3%", AMBER), ("Avg Time", "3.2 hr", "4hr SLA", NAVY)]
    sx2 = Inches(6.5)
    sy2 = Inches(2.28)
    for label, val, change, col in stats:
        card(s, sx2, sy2, Inches(1.4), Inches(0.75), fill=BG, border=TEAL_XL)
        rect(s, sx2, sy2, Inches(1.4), Inches(0.05), fill=col)
        tx(s, label, sx2 + Inches(0.08), sy2 + Inches(0.1),
           Inches(1.25), Inches(0.22), size=Pt(8.5), color=MUTED)
        tx(s, val, sx2 + Inches(0.08), sy2 + Inches(0.3),
           Inches(1.0), Inches(0.3), size=Pt(16), bold=True, color=NAVY)
        sx2 += Inches(1.5)

    # Chart area mock
    card(s, Inches(6.5), Inches(3.18), Inches(3.6), Inches(1.5), fill=BG, border=TEAL_XL)
    tx(s, "Messages — Last 7 Days", Inches(6.62), Inches(3.28),
       Inches(2.5), Inches(0.25), size=Pt(9), bold=True, color=NAVY)
    bars = [(0.45, TEAL_L), (0.65, TEAL_L), (0.5, TEAL_L),
            (0.85, TEAL), (0.6, TEAL_L), (0.95, TEAL), (0.7, TEAL_L)]
    days = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
    bx = Inches(6.75)
    for frac, col in bars:
        bh = frac * Inches(0.75)
        rect(s, bx, Inches(4.88) - bh, Inches(0.32), bh, fill=col)
        bx += Inches(0.38)

    # Message list mock
    card(s, Inches(10.25), Inches(3.18), Inches(2.5), Inches(1.5), fill=BG, border=TEAL_XL)
    tx(s, "Recent", Inches(10.35), Inches(3.28), Inches(2), Inches(0.25),
       size=Pt(9), bold=True, color=NAVY)
    msgs = [("alice@co.com", "Order #12345", "Auto", GREEN),
            ("+1 555 0000", "Refund request", "Esc", AMBER),
            ("bob@store.io", "Shipping status", "Auto", GREEN)]
    my2 = Inches(3.55)
    for sender, txt, status, col in msgs:
        tx(s, sender, Inches(10.35), my2, Inches(2.3), Inches(0.22),
           size=Pt(8.5), color=NAVY)
        tx(s, txt, Inches(10.35), my2 + Inches(0.2), Inches(1.7),
           Inches(0.2), size=Pt(8), color=MUTED)
        pill(s, Inches(11.7), my2 + Inches(0.2), status, bg=col)
        my2 += Inches(0.42)


# ── SLIDE 5: AI Pipeline ──────────────────────────────────────────────────────
def s5_ai():
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    header_bar(s, "How It Works: The AI Pipeline",
              "From inbound message to intelligent reply — in seconds, fully automated")
    footer(s)

    steps = [
        (TEAL,   "1", "Message Received",
         "WhatsApp or Email arrives via webhook.\n"
         "Metadata + body captured, timestamped,\nand routed to the AI pipeline."),
        (TEAL_L, "2", "Intent Classified",
         "LLM-powered classifier identifies the\n"
         "customer's intent: refund, order status,\n"
         "complaint, or general inquiry."),
        (NAVY,   "3", "Knowledge Retrieved",
         "RAG pipeline searches the Document Hub\n"
         "for relevant content — pricing, policies,\n"
         "product guides, past resolutions."),
        (TEAL_L, "4", "Decision Made",
         "AutoReplyService evaluates: should we\n"
         "auto-reply or escalate? Confidence score\n"
         "determines the routing path."),
        (TEAL,   "5", "Reply Sent",
         "AI-generated, contextually accurate reply\n"
         "sent via SendGrid or WhatsApp API.\n"
         "Full audit log written to database."),
    ]
    sw = Inches(2.2)
    sx = Inches(0.5)
    sy = Inches(1.9)
    for i, (col, num, title, desc) in enumerate(steps):
        card(s, sx, sy, sw, Inches(3.5), fill=col)
        # Number
        card(s, sx + Inches(0.8), sy + Inches(0.15),
             Inches(0.6), Inches(0.6), fill=WHITE)
        tx(s, num, sx + Inches(0.8), sy + Inches(0.17),
           Inches(0.6), Inches(0.6),
           size=Pt(16), bold=True, color=col, align=PP_ALIGN.CENTER)
        tx(s, title, sx, sy + Inches(0.88), sw, Inches(0.5),
           size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tx(s, desc, sx + Inches(0.12), sy + Inches(1.45),
           sw - Inches(0.24), Inches(1.9),
           size=Pt(11.5), color=WHITE, align=PP_ALIGN.CENTER, wrap=True)
        # Arrow
        if i < len(steps) - 1:
            tx(s, "→", sx + sw + Inches(0.05), sy + Inches(1.3),
               Inches(0.4), Inches(0.5),
               size=Pt(26), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        sx += sw + Inches(0.55)

    # Learning loop
    card(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.2),
         fill=WHITE, border=TEAL_XL)
    rect(s, Inches(0.5), Inches(5.6), Inches(0.08), Inches(1.2), fill=GREEN)
    tx(s, "Continuous Learning Loop", Inches(0.8), Inches(5.72),
       Inches(4), Inches(0.35), size=Pt(13), bold=True, color=NAVY)
    tx(s, "After every escalation is resolved, the LearningLoop records the operator's feedback — improving "
          "future classification accuracy and reply quality. Every resolved case makes the AI smarter for all customers.",
       Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.6),
       size=Pt(11.5), color=BODY, wrap=True)


# ── SLIDE 6: Channels ─────────────────────────────────────────────────────────
def s6_channels():
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    header_bar(s, "Multi-Channel: WhatsApp + Email",
              "Both channels share the same AI pipeline, knowledge base, and operator portal")
    footer(s)

    # WhatsApp card
    card(s, Inches(0.65), Inches(1.75), Inches(5.7), Inches(4.6),
         fill=WHITE, border=TEAL_XL)
    rect(s, Inches(0.65), Inches(1.75), Inches(5.7), Inches(0.65), fill=NAVY)
    tx(s, "WhatsApp", Inches(0.9), Inches(1.9), Inches(4), Inches(0.45),
       size=Pt(18), bold=True, color=WHITE)
    wa = [
        "Meta Cloud API — modern, stable, officially supported",
        "Webhook verification with HMAC-SHA256 signatures",
        "Template messages for outbound notifications",
        "Contact profiling and phone number validation",
        "Media: images, documents, audio, voice notes",
        "Delivery and read receipts tracked in portal",
        "Opt-out / block handling fully automated",
        "Business profile management via API",
    ]
    wy = Inches(2.6)
    for f in wa:
        tx(s, f"  {f}", Inches(0.85), wy, Inches(5.3), Inches(0.35),
           size=Pt(12), color=BODY)
        wy += Inches(0.46)

    # Email card
    card(s, Inches(6.7), Inches(1.75), Inches(6.1), Inches(4.6),
         fill=WHITE, border=TEAL_XL)
    rect(s, Inches(6.7), Inches(1.75), Inches(6.1), Inches(0.65), fill=TEAL)
    tx(s, "Email", Inches(6.95), Inches(1.9), Inches(4), Inches(0.45),
       size=Pt(18), bold=True, color=WHITE)
    em = [
        "SendGrid Inbound Parse — full email receipt with attachments",
        "Attachment ingestion: files saved to Document Hub automatically",
        "HTML and plain-text body parsing, intelligently normalized",
        "Reply-to threading via In-Reply-To and References headers",
        "SPF / DKIM / DMARC aware routing and validation",
        "Auto-reply with quoted original for customer clarity",
        "PDPA: email content minimized in storage post-resolution",
        "Bounce handling and suppression list management",
    ]
    ey = Inches(2.6)
    for f in em:
        tx(s, f"  {f}", Inches(6.9), ey, Inches(5.7), Inches(0.35),
           size=Pt(12), color=BODY)
        ey += Inches(0.46)

    # Bottom note
    tx(s, "Both channels: same AI pipeline — classify → RAG → decide → reply",
       Inches(0.65), Inches(6.55), Inches(12.1), Inches(0.35),
       size=Pt(12), bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ── SLIDE 7: Escalations ───────────────────────────────────────────────────────
def s7_escalations():
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    header_bar(s, "Smart Escalations",
              "The right message, to the right operator, with full context — instantly")
    footer(s)

    # Left: triggers
    tx(s, "Escalation Triggers", Inches(0.65), Inches(1.7), Inches(5),
       size=Pt(15), bold=True, color=NAVY)
    triggers = [
        (RED,   "Cancellation / Refund Threats",
         "High-friction intent signals — classifier confidence > 0.8 triggers immediate escalation."),
        (AMBER, "VIP / Priority Contacts",
         "Contacts flagged by tier or manually tagged receive instant priority routing."),
        (AMBER, "SLA Breach Risk",
         "Messages unresolved past the configured SLA window auto-escalate."),
        (RED,   "Negative Sentiment",
         "Frustration or anger detected by the classifier — regardless of content topic."),
        (RED,   "Repeat Queries (3x+)",
         "Same topic from same contact, unanswered, escalates after 3 attempts."),
        (TEAL,  "Manual Flag",
         "Operator can flag any thread for escalation with one click and a note."),
    ]
    ty = Inches(2.1)
    for col, title, desc in triggers:
        rect(s, Inches(0.65), ty + Inches(0.06), Inches(0.06), Inches(0.4), fill=col)
        tx(s, title, Inches(0.85), ty, Inches(2.5), Inches(0.28),
           size=Pt(12), bold=True, color=NAVY)
        tx(s, desc, Inches(0.85), ty + Inches(0.28), Inches(5.3),
           Inches(0.35), size=Pt(10.5), color=BODY, wrap=True)
        ty += Inches(0.72)

    # Right: flow
    card(s, Inches(6.8), Inches(1.65), Inches(6.1), Inches(5.1),
         fill=WHITE, border=TEAL_XL)
    rect(s, Inches(6.8), Inches(1.65), Inches(6.1), Inches(0.55), fill=NAVY)
    tx(s, "Resolution Flow", Inches(7.05), Inches(1.75), Inches(5.5),
       size=Pt(14), bold=True, color=WHITE)

    flow = [
        ("1", "Alert Created",
         "Escalation record created with full message context, contact info, and SLA countdown."),
        ("2", "Assigned",
         "Auto-assigned to available operator based on team routing rules and contact tier."),
        ("3", "Context Delivered",
         "Portal shows full thread + relevant document excerpts from RAG — no context switching."),
        ("4", "Operator Resolves",
         "Outcome logged (resolved / transferred / closed). LearningLoop records for training."),
        ("5", "SLA Closed",
         "Ticket closed. Resolution time, outcome, and satisfaction stored in immutable audit trail."),
    ]
    fy = Inches(2.4)
    for num, title, desc in flow:
        card(s, Inches(7.0), fy, Inches(0.5), Inches(0.5), fill=TEAL)
        tx(s, num, Inches(7.0), fy + Inches(0.04), Inches(0.5), Inches(0.45),
           size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tx(s, title, Inches(7.65), fy + Inches(0.04), Inches(2.5), Inches(0.28),
           size=Pt(12), bold=True, color=NAVY)
        tx(s, desc, Inches(7.65), fy + Inches(0.32), Inches(5.0),
           Inches(0.42), size=Pt(10.5), color=BODY, wrap=True)
        if num != "5":
            rect(s, Inches(7.22), fy + Inches(0.54), Inches(0.06), Inches(0.2), fill=TEAL_XL)
        fy += Inches(0.82)


# ── SLIDE 8: PDPA ─────────────────────────────────────────────────────────────
def s8_compliance():
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    header_bar(s, "PDPA Compliance by Design",
              "Every layer of Sequor is built with Singapore data protection requirements at its core")
    footer(s)

    pillars = [
        (TEAL,   "1", "Data Minimization",
         "Sequor stores only operationally necessary data. Email bodies and WhatsApp content are "
         "retained only for active threads. Historical archives are pseudonymized."),
        (NAVY,   "2", "Right to Erasure",
         "Full deletion workflow: data subject requests erasure → all PII purged from active records "
         "and audit logs. Erasure confirmation issued automatically."),
        (TEAL_L, "3", "Immutable Audit Trail",
         "Every message, classification, escalation, and operator action logged with tenant_id, timestamp, "
         "actor, and outcome — PDPA Art. 20 compliant."),
        (TEAL,   "4", "Tenant Isolation",
         "Role-based access: operators see only their assigned contacts. Database-level tenant isolation "
         "enforced on every query. No cross-tenant data leakage."),
        (NAVY,   "5", "Consent Management",
         "All inbound contacts flagged for consent status. Audit log records channel, notice version, "
         "and timestamp of consent — retrievable on demand."),
        (TEAL_L, "6", "PII Classification",
         "Email addresses, phone numbers, and names classified as PII in the audit trail. RAG retrieval "
         "respects classification boundaries — no PII in AI training data."),
    ]
    pw = Inches(3.8)
    px = Inches(0.65)
    py = Inches(1.75)
    for i, (col, num, title, desc) in enumerate(pillars):
        row = i // 3
        col_i = i % 3
        x = px + col_i * (pw + Inches(0.28))
        y = py + row * Inches(2.55)
        card(s, x, y, pw, Inches(2.35), fill=WHITE, border=TEAL_XL)
        rect(s, x, y, pw, Inches(0.08), fill=col)
        card(s, x + Inches(0.2), y + Inches(0.18), Inches(0.5), Inches(0.5), fill=col)
        tx(s, num, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5),
           size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tx(s, title, x + Inches(0.85), y + Inches(0.25), pw - Inches(1.0),
           Inches(0.35), size=Pt(13), bold=True, color=NAVY)
        tx(s, desc, x + Inches(0.22), y + Inches(0.78), pw - Inches(0.4),
           Inches(1.45), size=Pt(11), color=BODY, wrap=True)


# ── SLIDE 9: Pricing ───────────────────────────────────────────────────────────
def s9_pricing():
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    header_bar(s, "Pricing",
              "Simple, transparent, scalable — free forever for development and testing")
    footer(s)

    plans = [
        ("Free",    "$0",    "month",
         ["50 messages / month", "WhatsApp + Email channels",
          "Auto-reply AI", "RAG Document Hub",
          "Dashboard & message history",
          "1 operator account",
          "Community support"],
         TEAL_XL, NAVY, False),
        ("Solo",    "$15",   "month",
         ["200 messages / month", "WhatsApp + Email channels",
          "Auto-reply AI + RAG",
          "Dashboard & message history",
          "30-day message retention",
          "1 operator account",
          "Email support"],
         TEAL_XL, NAVY, False),
        ("Starter", "$35",   "month",
         ["Unlimited messages*", "WhatsApp + Email channels",
          "Auto-reply AI + RAG",
          "Smart Escalations",
          "PDPA audit trail",
          "90-day message retention",
          "Up to 3 operators",
          "Priority support"],
         TEAL, WHITE, True),
        ("Pro",     "$55",   "month",
         ["Unlimited messages", "Everything in Starter",
          "Unlimited operators",
          "Unlimited documents",
          "Priority RAG processing",
          "Advanced analytics dashboard",
          "Custom key phrase maps",
          "PDPA compliance report",
          "Dedicated support"],
         NAVY, WHITE, False),
    ]

    pw = Inches(2.9)
    px = Inches(0.55)
    py = Inches(1.8)
    for name, price, period, features, col, txt, featured in plans:
        card_h = Inches(5.1)
        card(s, px, py, pw, card_h, fill=col,
             border=col if featured else TEAL_XL,
             border_w=Pt(2) if featured else Pt(1))
        if featured:
            rect(s, px, py, pw, Inches(0.42), fill=TEAL_L)
            tx(s, "Most Popular", px, py + Inches(0.08), pw, Inches(0.32),
               size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            ty_name = py + Inches(0.55)
        else:
            ty_name = py + Inches(0.22)
        tx(s, name, px, ty_name, pw, Inches(0.55),
           size=Pt(22), bold=True, color=txt, align=PP_ALIGN.CENTER)
        tx(s, price, px, ty_name + Inches(0.6), pw, Inches(0.75),
           size=Pt(40), bold=True, color=txt, align=PP_ALIGN.CENTER)
        tx(s, f"/{period}", px, ty_name + Inches(1.3), pw, Inches(0.3),
           size=Pt(12), color=txt, align=PP_ALIGN.CENTER)
        add_line(s, px + Inches(0.3), ty_name + Inches(1.7),
                 pw - Inches(0.6), Inches(0.02), color=txt if featured else TEAL_XL)
        fy = ty_name + Inches(1.85)
        for feat in features:
            tx(s, f"✓  {feat}", px + Inches(0.25), fy,
               pw - Inches(0.4), Inches(0.36), size=Pt(11), color=txt)
            fy += Inches(0.4)
        px += pw + Inches(0.22)

    tx(s, "* Message fair-use limits enforced at 3x plan average. Enterprise volume pricing available on request.",
       Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.3),
       size=Pt(10), color=MUTED, italic=True)


# ── SLIDE 10: Setup ───────────────────────────────────────────────────────────
def s10_setup():
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    header_bar(s, "Getting Started in Minutes",
              "No servers, no ML pipelines, no data science team required")
    footer(s)

    # Left: timeline
    steps = [
        (TEAL,   "1", "Connect WhatsApp Business Account",
         "Log in to Meta Business Manager → add the WhatsApp Business API — "
         "copy your phone number ID and access token into the Sequor portal. Takes 3 minutes."),
        (NAVY,   "2", "Connect Email via SendGrid",
         "Create a free SendGrid account → configure Inbound Parse → "
         "point the webhook at your Sequor URL. Copy your API key in. Done."),
        (TEAL_L, "3", "Upload Your Documents",
         "Drop your PDFs, policy docs, FAQs, and product guides into the Document Hub. "
         "RAG indexing starts automatically — takes ~5 minutes for 50 documents."),
        (TEAL,   "4", "Configure Key Phrases & Escalation Rules",
         "Add the topics your customers ask about. Map them to reply strategies. "
         "Set your SLA window and escalation contacts. 10 minutes."),
        (NAVY,   "5", "Go Live",
         "Flip the channel to Active. Every inbound message now routes through the AI pipeline. "
         "First auto-reply typically fires within the hour."),
    ]
    sw = Inches(5.8)
    sx = Inches(0.65)
    sy = Inches(1.75)
    for col, num, title, desc in steps:
        card(s, sx, sy, Inches(0.55), Inches(0.9), fill=col)
        tx(s, num, sx, sy + Inches(0.04), Inches(0.55), Inches(0.9),
           size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tx(s, title, sx + Inches(0.7), sy, sw - Inches(0.7), Inches(0.35),
           size=Pt(13), bold=True, color=NAVY)
        tx(s, desc, sx + Inches(0.7), sy + Inches(0.38), sw - Inches(0.7),
           Inches(0.5), size=Pt(11), color=BODY, wrap=True)
        sy += Inches(1.02)

    # Right: what RAG requires
    card(s, Inches(7.0), Inches(1.75), Inches(5.9), Inches(4.9),
         fill=WHITE, border=TEAL_XL)
    rect(s, Inches(7.0), Inches(1.75), Inches(5.9), Inches(0.6), fill=NAVY)
    tx(s, "About Document Quality", Inches(7.25), Inches(1.9),
       Inches(5.5), Inches(0.4), size=Pt(15), bold=True, color=WHITE)
    tx(s,
       "The AI's accuracy is directly proportional to your document quality. "
       "Well-structured PDFs with clear headings, Q&A format, or policy documents "
       "produce the best auto-replies. Poorly formatted scans perform accordingly.",
       Inches(7.25), Inches(2.5), Inches(5.4), Inches(1.0),
       size=Pt(12), color=BODY, wrap=True)
    add_line(s, Inches(7.25), Inches(3.6), Inches(5.4))
    tx(s, "Best document types:", Inches(7.25), Inches(3.72),
       Inches(5.4), Inches(0.3), size=Pt(12), bold=True, color=NAVY)
    best = [
        "Product FAQ sheets (Q&A format)",
        "Policy documents (pricing, shipping, returns)",
        "Service descriptions and menus",
        "How-to guides and troubleshooting docs",
        "Pricing sheets and plan comparisons",
    ]
    by = Inches(4.05)
    for b in best:
        tx(s, f"✓  {b}", Inches(7.25), by, Inches(5.4), Inches(0.33),
           size=Pt(11.5), color=BODY)
        by += Inches(0.37)


# ── SLIDE 11: Live Demo ────────────────────────────────────────────────────────
def s11_demo():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, fill=DARK)
    rect(s, Inches(8.5), 0, Inches(4.83), H, fill=NAVY)
    rect(s, Inches(9.8), 0, Inches(3.53), H, fill=RGBColor(0x1D, 0x5A, 0x77))

    tx(s, "Live Demo", Inches(0.8), Inches(1.0), Inches(7), Inches(0.9),
       size=Pt(44), bold=True, color=WHITE)
    add_line(s, Inches(0.8), Inches(2.0), Inches(4), color=TEAL)
    tx(s, "Watch Sequor handle a real inbound message end-to-end",
       Inches(0.8), Inches(2.15), Inches(7.5), Inches(0.5),
       size=Pt(16), color=TEAL_L)

    demo_steps = [
        ("1", "Inbound message arrives", "via WhatsApp or Email webhook"),
        ("2", "AI classifies intent", "and retrieves relevant knowledge from Document Hub"),
        ("3", "Auto-reply generated", "and sent — or escalation triggered to operator"),
        ("4", "Portal updates", "dashboard, message history, and escalation queue"),
        ("5", "Escalation resolved", "operator closes ticket, AI learns for next time"),
    ]
    dx = Inches(0.8)
    dy = Inches(3.1)
    for num, step, detail in demo_steps:
        card(s, dx, dy, Inches(0.55), Inches(0.55), fill=TEAL)
        tx(s, num, dx, dy + Inches(0.04), Inches(0.55), Inches(0.55),
           size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tx(s, step, dx + Inches(0.72), dy + Inches(0.04), Inches(3), Inches(0.32),
           size=Pt(14), bold=True, color=WHITE)
        tx(s, detail, dx + Inches(0.72), dy + Inches(0.34), Inches(6), Inches(0.28),
           size=Pt(12), color=TEAL_L)
        dy += Inches(0.74)

    tx(s, "Try it now → portal.sequor.com",
       Inches(0.8), Inches(6.6), Inches(7), Inches(0.4),
       size=Pt(14), bold=True, color=TEAL)

    tx(s, "Portal\nLive",
       Inches(9.0), Inches(2.5), Inches(3.5), Inches(1.5),
       size=Pt(26), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, "Sequor Pte. Ltd.\nSingapore",
       Inches(9.0), Inches(4.1), Inches(3.5), Inches(0.7),
       size=Pt(12), color=TEAL_L, align=PP_ALIGN.CENTER)


# ── SLIDE 12: Roadmap ──────────────────────────────────────────────────────────
def s12_roadmap():
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    header_bar(s, "Roadmap",
              "What's coming — building the full customer communication platform")
    footer(s)

    quarters = [
        ("Q3 2026", TEAL, [
            ("SMS Channel",
             "Twilio integration for SMS inbound/outbound — complete omnichannel coverage."),
            ("Team Roles & Permissions",
             "Admin, Manager, Operator tiers with granular access controls."),
            ("Analytics Dashboard",
             "Message volume trends, AI accuracy reports, SLA performance tracking."),
        ]),
        ("Q4 2026", NAVY, [
            ("Multi-Language AI",
             "Thai, Bahasa, Vietnamese, Mandarin — AI replies in the customer's language."),
            ("API Webhooks",
             "Custom outbound triggers for CRM and ERP integration (Salesforce, HubSpot)."),
            ("Mobile App",
             "iOS and Android — escalation management and approvals on the go."),
        ]),
        ("Q1 2027", RGBColor(0x1D, 0x5A, 0x77), [
            ("Enterprise SSO",
             "SAML and OIDC single sign-on for large enterprise procurement."),
            ("SOC 2 Type II",
             "Security certification — required for Fortune 500 and regulated industries."),
            ("Voice Channel",
             "Inbound voice with transcription and AI summarization before routing."),
        ]),
    ]
    rw = Inches(3.9)
    rx = Inches(0.65)
    ry = Inches(1.8)
    for qtr, col, items in quarters:
        card(s, rx, ry, rw, Inches(5.0), fill=WHITE, border=TEAL_XL)
        rect(s, rx, ry, rw, Inches(0.65), fill=col)
        tx(s, qtr, rx, ry + Inches(0.13), rw, Inches(0.45),
           size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        iy = ry + Inches(0.85)
        for title, desc in items:
            rect(s, rx + Inches(0.22), iy + Inches(0.08),
                 Inches(0.06), Inches(0.4), fill=col)
            tx(s, title, rx + Inches(0.42), iy, rw - Inches(0.6),
               Inches(0.35), size=Pt(13), bold=True, color=NAVY)
            tx(s, desc, rx + Inches(0.42), iy + Inches(0.38),
               rw - Inches(0.6), Inches(0.6), size=Pt(11.5), color=BODY, wrap=True)
            iy += Inches(1.2)
        rx += rw + Inches(0.28)


# ── SLIDE 13: Why Sequor ──────────────────────────────────────────────────────
def s13_why():
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    header_bar(s, "Why Sequor",
              "Built for Southeast Asian SMBs who need enterprise-grade AI without the complexity")
    footer(s)

    reasons = [
        (TEAL,   "Designed for Southeast Asia",
         "WhatsApp as the primary channel, multi-language AI, PDPA-first compliance — "
         "designed for the region's SMB reality, not adapted from Western enterprise software."),
        (NAVY,   "10x Operator Efficiency",
         "80% of messages auto-replied. Operators spend time on the 20% that matters — "
         "escalations, VIPs, and complex cases — not repetitive FAQs."),
        (TEAL_L, "Zero Infrastructure Headaches",
         "No servers, no ML pipelines, no training data to build. Connect your channels "
         "and upload your documents. The AI handles the rest."),
        (TEAL,   "PDPA from Day One",
         "Compliance isn't an afterthought or an add-on module. Audit trails, erasure "
         "workflows, and data minimization are built into every layer."),
        (NAVY,   "Transparent Pricing",
         "Predictable per-seat pricing. No per-message AI credits, no hidden ingestion "
         "fees, no enterprise-only features locked behind a sales call."),
        (TEAL_L, "Fast Time to Value",
         "Channels connected in minutes. First AI reply fires within the hour. "
         "Document Hub built from your existing files — no structured data required."),
    ]
    rw = Inches(3.9)
    rx = Inches(0.65)
    ry = Inches(1.8)
    for i, (col, title, desc) in enumerate(reasons):
        row = i // 3
        col_i = i % 3
        x = rx + col_i * (rw + Inches(0.28))
        y = ry + row * Inches(2.5)
        card(s, x, y, rw, Inches(2.2), fill=WHITE, border=TEAL_XL)
        rect(s, x, y, rw, Inches(0.08), fill=col)
        tx(s, title, x + Inches(0.22), y + Inches(0.2), rw - Inches(0.4),
           Inches(0.5), size=Pt(14), bold=True, color=NAVY)
        tx(s, desc, x + Inches(0.22), y + Inches(0.75), rw - Inches(0.4),
           Inches(1.3), size=Pt(12), color=BODY, wrap=True)


# ── SLIDE 14: Opportunity ────────────────────────────────────────────────────
def s14_opportunity():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, fill=DARK)
    rect(s, Inches(8.5), 0, Inches(4.83), H, fill=NAVY)
    rect(s, Inches(9.8), 0, Inches(3.53), H, fill=RGBColor(0x1D, 0x5A, 0x77))

    tx(s, "The Opportunity", Inches(0.8), Inches(0.7), Inches(7.5), Inches(0.9),
       size=Pt(38), bold=True, color=WHITE)
    add_line(s, Inches(0.8), Inches(1.7), Inches(4), color=TEAL)

    stats = [
        ("$2.4B", "TAM — SE Asian SMB Customer\nCommunication Software (2026)"),
        ("35%", "Annual growth in WhatsApp\nBusiness API adoption"),
        ("80%", "Average auto-reply rate\nin Sequor pilot accounts"),
        ("< 4 wks", "Customer time-to-first\nAI-replied message"),
    ]
    sx = Inches(0.8)
    sy = Inches(2.1)
    for val, desc in stats:
        tx(s, val, sx, sy, Inches(2.6), Inches(0.85),
           size=Pt(36), bold=True, color=TEAL)
        tx(s, desc, sx, sy + Inches(0.88), Inches(2.6), Inches(0.75),
           size=Pt(11.5), color=TEAL_L, wrap=True)
        sy += Inches(1.85)

    # Investment box
    card(s, Inches(8.5), Inches(1.65), Inches(4.5), Inches(5.1), fill=NAVY)
    rect(s, Inches(8.5), Inches(1.65), Inches(4.5), Inches(0.65), fill=TEAL)
    tx(s, "Investment Summary", Inches(8.7), Inches(1.78), Inches(4.1), Inches(0.45),
       size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("Raise", "Seed Round — S$1.2M"),
        ("Use of Funds",
         "Engineering (40%)\nSales & Marketing (35%)\nOperations (25%)"),
        ("Key Milestones",
         "20 paying customers\n1,000 monthly messages\n3 channel integrations"),
        ("Target Close", "Q3 2026"),
    ]
    ry2 = Inches(2.55)
    for label, val in rows:
        tx(s, label + ":", Inches(8.7), ry2, Inches(1.5), Inches(0.32),
           size=Pt(11), bold=True, color=TEAL_L)
        tx(s, val, Inches(10.25), ry2, Inches(2.6), Inches(0.65),
           size=Pt(11.5), color=WHITE, wrap=True)
        ry2 += Inches(0.82)

    footer(s)


# ── SLIDE 15: Contact ─────────────────────────────────────────────────────────
def s15_contact():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, fill=DARK)
    rect(s, Inches(8.5), 0, Inches(4.83), H, fill=NAVY)
    rect(s, Inches(9.8), 0, Inches(3.53), H, fill=RGBColor(0x1D, 0x5A, 0x77))

    # Logo
    card(s, Inches(0.8), Inches(1.5), Inches(0.85), Inches(0.85), fill=TEAL)
    tx(s, "S", Inches(0.8), Inches(1.52), Inches(0.85), Inches(0.85),
       size=Pt(38), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, "Sequor", Inches(1.85), Inches(1.55), Inches(4), Inches(0.8),
       size=Pt(38), bold=True, color=WHITE)
    add_line(s, Inches(0.8), Inches(2.55), Inches(5), color=TEAL)
    tx(s, "AI-Powered Customer Communication Platform",
       Inches(0.8), Inches(2.68), Inches(7), Inches(0.4),
       size=Pt(14), color=TEAL_L)

    contact_items = [
        ("🌐", "Website",  "www.sequor.com"),
        ("📧", "Email",    "hello@sequor.com"),
        ("💬", "WhatsApp", "+65 9000 0000"),
        ("📍", "Address",  "1 Raffles Place, #12-00\nSingapore 048616"),
    ]
    cy = Inches(3.3)
    for icon, label, val in contact_items:
        tx(s, f"{icon}  {label}:", Inches(0.8), cy, Inches(2.2), Inches(0.35),
           size=Pt(13), bold=True, color=TEAL_L)
        tx(s, val, Inches(3.0), cy, Inches(5), Inches(0.45),
           size=Pt(13), color=WHITE)
        cy += Inches(0.6)

    tx(s, "Thank you for\nyour time.",
       Inches(9.0), Inches(2.5), Inches(3.5), Inches(1.5),
       size=Pt(30), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_line(s, Inches(9.0), Inches(4.1), Inches(3.5), color=TEAL)
    tx(s, "We look forward to building\nthe future of customer communication.",
       Inches(9.0), Inches(4.3), Inches(3.5), Inches(0.9),
       size=Pt(12), color=TEAL_L, align=PP_ALIGN.CENTER, wrap=True)

    footer(s)


# ── Build ─────────────────────────────────────────────────────────────────────
s1_cover()
s2_problem()
s3_solution()
s4_product()
s5_ai()
s6_channels()
s7_escalations()
s8_compliance()
s9_pricing()
s10_setup()
s11_demo()
s12_roadmap()
s13_why()
s14_opportunity()
s15_contact()

out = "/Users/aliciapang/Documents/GitHub/Sequor/workspaces/_template/Sequor_Investor_Deck_May2026.pptx"
prs.save(out)
print(f"Saved: {out}")
