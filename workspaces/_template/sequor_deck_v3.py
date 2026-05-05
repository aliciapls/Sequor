"""Generate Sequor PE Investor Presentation Deck — Professional Version v3."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt as PT

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1A, 0x4A, 0x63)
TEAL    = RGBColor(0x3C, 0x8E, 0xAF)
TEAL_L  = RGBColor(0x5F, 0xA6, 0xC5)
TEAL_XL = RGBColor(0xB7, 0xDD, 0xEA)
BG      = RGBColor(0xF0, 0xF9, 0xFB)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BODY    = RGBColor(0x47, 0x56, 0x6A)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
RED     = RGBColor(0xEF, 0x44, 0x44)
DARK    = RGBColor(0x0F, 0x17, 0x2A)
NAVY2   = RGBColor(0x1D, 0x5A, 0x77)

# Font: Calibri is guaranteed on all Macs/Windows — use it throughout
FONT = "Calibri"

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Core helpers ────────────────────────────────────────────────────────────────

def mkrect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.line.width = lw
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s


def mktx(slide, text, l, t, w, h,
          size=18, bold=False, italic=False,
          color=BODY, align=PP_ALIGN.LEFT, wrap=True):
    """Create a textbox with properly formatted text."""
    from pptx.util import Pt
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    # Vertical centering hint
    from pptx.enum.text import MSO_ANCHOR
    try:
        tf.auto_size = None
    except Exception:
        pass
    return tb


def mktx_vcenter(slide, text, l, t, w, h,
                  size=18, bold=False, italic=False,
                  color=BODY, align=PP_ALIGN.LEFT):
    """Textbox with vertically centered text via paragraph spacing."""
    from pptx.util import Pt
    from pptx.enum.text import MSO_ANCHOR
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    # Center vertically by setting line spacing to match box height
    from pptx.util import Pt as PT
    from pptx.oxml.ns import qn
    pPr = p._p.get_or_add_pPr()
    pPr.set('algn', 'ctr' if align == PP_ALIGN.CENTER else 'l')
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    try:
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    return tb


def mktx_lines(slide, lines, l, t, w, line_h,
               size=13, color=BODY, bold_first=False):
    """Multi-line text box. lines = list of (text, bold, color)."""
    from pptx.util import Pt
    from pptx.enum.text import MSO_ANCHOR
    # Calculate total height
    n = len(lines)
    total_h = n * line_h + Pt(12)
    tb = slide.shapes.add_textbox(l, t, w, total_h)
    tf = tb.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.TOP
    except Exception:
        pass
    for i, (text, bold, c) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = c
        r.font.name = FONT
    return tb


def hline(slide, l, t, w, color=TEAL, thick=Pt(0.8)):
    mkrect(slide, l, t, w, thick, fill=color)


def footer_bar(slide, text="Confidential  •  May 2026  •  Sequor"):
    mkrect(slide, 0, Inches(7.1), W, Inches(0.4), fill=NAVY)
    mktx(slide, text, Inches(0.65), Inches(7.13), Inches(8), Inches(0.32),
          size=10, color=TEAL_L, italic=True)


def header(slide, title, subtitle=None):
    mkrect(slide, 0, 0, W, Inches(1.35), fill=NAVY)
    mktx(slide, title, Inches(0.65), Inches(0.28), Inches(10), Inches(0.6),
          size=30, bold=True, color=WHITE)
    if subtitle:
        mktx(slide, subtitle, Inches(0.65), Inches(0.88), Inches(11), Inches(0.38),
              size=13, color=TEAL_L, italic=True)


def bg(slide):
    mkrect(slide, 0, 0, W, H, fill=BG)


def pill(slide, l, t, text, bg_col=TEAL, fg=WHITE, w=Inches(0.85)):
    mkrect(slide, l, t, w, Inches(0.3), fill=bg_col)
    mktx(slide, text, l, t + Inches(0.04), w, Inches(0.25),
          size=9, bold=True, color=fg, align=PP_ALIGN.CENTER)


# ── SLIDE 1: Cover ─────────────────────────────────────────────────────────────
def s1():
    s = prs.slides.add_slide(BLANK)
    mkrect(s, 0, 0, W, H, fill=DARK)
    mkrect(s, Inches(8.2), 0, Inches(5.13), H, fill=NAVY)
    mkrect(s, Inches(9.7), 0, Inches(3.63), H, fill=NAVY2)

    # Logo
    mkrect(s, Inches(0.75), Inches(1.4), Inches(0.85), Inches(0.85), fill=TEAL)
    mktx(s, "S", Inches(0.75), Inches(1.42), Inches(0.85), Inches(0.85),
          size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    mktx(s, "Sequor", Inches(1.78), Inches(1.47), Inches(4), Inches(0.75),
          size=38, bold=True, color=WHITE)

    hline(s, Inches(0.75), Inches(2.4), Inches(6.5))

    mktx(s, "AI-Powered Customer Communication Platform",
          Inches(0.75), Inches(2.52), Inches(7), Inches(0.42),
          size=15, color=TEAL_L)

    mktx(s, "The operating system for\ncustomer conversations.",
          Inches(0.75), Inches(3.2), Inches(7), Inches(1.6),
          size=38, bold=True, color=WHITE)

    mktx(s, "One platform. Every channel. Fully automated.\nBuilt for Southeast Asian SMBs — PDPA compliant from day one.",
          Inches(0.75), Inches(5.1), Inches(7), Inches(0.8),
          size=13, color=MUTED)

    # Right panel
    mktx(s, "Investor\nPresentation",
          Inches(9.0), Inches(2.3), Inches(3.5), Inches(1.4),
          size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    mktx(s, "May 2026",
          Inches(9.0), Inches(3.85), Inches(3.5), Inches(0.38),
          size=14, color=TEAL_L, align=PP_ALIGN.CENTER)
    mktx(s, "Sequor Pte. Ltd.\nSingapore",
          Inches(9.0), Inches(5.4), Inches(3.5), Inches(0.7),
          size=12, color=MUTED, align=PP_ALIGN.CENTER)

    footer_bar(s)


# ── SLIDE 2: Problem ─────────────────────────────────────────────────────────
def s2():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, "The Problem",
           "Why businesses are losing customers — and operators — every day")
    footer_bar(s)

    problems = [
        (RED,   "Channels in silos",
         "WhatsApp, Email, SMS — each a separate inbox. Context lost at every handoff."),
        (AMBER, "Slow responses kill loyalty",
         "Average first response: 4–6 hours. 35% of customers leave after one slow reply."),
        (RED,   "Repetitive queries burn out teams",
         "70% of inbound messages are repeats — FAQs, order status, product info."),
        (AMBER, "Escalations slip through",
         "High-intent signals missed. Resolving one takes 6+ tools and multiple tabs."),
        (RED,   "Compliance is manual and risky",
         "PDPA requires audit trails, erasure, data minimization. Manual processes fail audits."),
        (MUTED, "Zero operational visibility",
         "No unified view of volume, resolution rates, or AI accuracy."),
    ]

    cw = Inches(3.85)
    ch = Inches(1.62)
    sx = Inches(0.62)
    sy = Inches(1.65)
    gx = Inches(0.25)
    gy = Inches(0.2)

    for i, (col, title, desc) in enumerate(problems):
        col_i = i % 3
        row = i // 3
        x = sx + col_i * (cw + gx)
        y = sy + row * (ch + gy)
        mkrect(s, x, y, cw, ch, fill=WHITE, line=TEAL_XL, lw=Pt(0.75))
        mkrect(s, x, y, Inches(0.07), ch, fill=col)
        mktx(s, title, x + Inches(0.2), y + Inches(0.18), cw - Inches(0.35), Inches(0.32),
              size=13, bold=True, color=NAVY)
        mktx(s, desc, x + Inches(0.2), y + Inches(0.55), cw - Inches(0.35), Inches(0.95),
              size=11.5, color=BODY, wrap=True)


# ── SLIDE 3: Solution ───────────────────────────────────────────────────────
def s3():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, "The Solution",
           "One platform — every channel, every message, fully automated")
    footer_bar(s)

    mkrect(s, Inches(0.62), Inches(1.68), Inches(8.0), Inches(1.65),
           fill=WHITE, line=TEAL_XL, lw=Pt(0.75))
    mktx(s,
          "Sequor unifies WhatsApp and Email into a single AI-powered platform — auto-replies to "
          "common queries, intelligently escalates what matters, and gives operators complete "
          "visibility — while remaining PDPA compliant by design.",
          Inches(0.85), Inches(1.82), Inches(7.6), Inches(1.4),
          size=13.5, color=BODY, wrap=True)

    pillars = [
        (TEAL,   "Auto-Reply AI",
         "Classifies intent, retrieves relevant documents via RAG, and generates accurate replies — without human input."),
        (NAVY,   "Smart Escalations",
         "Detects high-intent signals and routes to the right operator with full conversation context."),
        (TEAL_L, "PDPA Compliance",
         "Built-in audit trail, data minimization, and erasure workflows — compliant from day one."),
    ]
    pw = Inches(3.85)
    px = Inches(0.62)
    py = Inches(3.6)
    for col, title, desc in pillars:
        mkrect(s, px, py, pw, Inches(1.72), fill=WHITE, line=TEAL_XL, lw=Pt(0.75))
        mkrect(s, px, py, pw, Inches(0.07), fill=col)
        mktx(s, title, px + Inches(0.22), py + Inches(0.2), pw - Inches(0.4), Inches(0.38),
              size=15, bold=True, color=NAVY)
        mktx(s, desc, px + Inches(0.22), py + Inches(0.65), pw - Inches(0.4), Inches(0.95),
              size=12, color=BODY, wrap=True)
        px += pw + Inches(0.28)

    metrics = [("80%", "messages auto-replied"), ("4-hr", "avg SLA achievement"),
                ("< 10 min", "to first AI reply"), ("0", "servers to manage")]
    mx = Inches(0.62)
    my = Inches(5.58)
    for val, lbl in metrics:
        mktx(s, val, mx, my, Inches(2.8), Inches(0.62),
              size=34, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        mktx(s, lbl, mx, my + Inches(0.65), Inches(2.8), Inches(0.35),
              size=11, color=MUTED, align=PP_ALIGN.CENTER)
        mx += Inches(3.1)


# ── SLIDE 4: Product ──────────────────────────────────────────────────────────
def s4():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, "Product: The Operator Portal",
           "Your command center for every customer interaction — one view, all channels")
    footer_bar(s)

    features = [
        (TEAL,   "Dashboard",
         "Live metrics: message volume, auto-reply rate, SLA health, and escalation queue at a glance."),
        (NAVY,   "Messages",
         "Unified thread history across WhatsApp and Email — full context without switching tabs."),
        (TEAL,   "Escalations",
         "Intelligent queue ranked by priority — assign, resolve, and track resolution time."),
        (NAVY,   "Auto-Replies",
         "Log of every AI-generated reply with confidence score and operator feedback loop."),
        (TEAL,   "Document Hub",
         "Knowledge base linked to AI RAG — upload documents, AI uses them in every reply."),
        (NAVY,   "Key Phrases",
         "Custom trigger phrases mapped to AI reply strategies per topic or channel."),
        (TEAL,   "Channels",
         "WhatsApp and Email connections — webhook status, health, and configuration."),
        (NAVY,   "Subscription",
         "Plan, usage, and billing — transparent, predictable, self-serve."),
    ]
    fx = Inches(0.62)
    fy = Inches(1.68)
    fw = Inches(5.2)
    for col, title, desc in features:
        mkrect(s, fx, fy + Inches(0.06), Inches(0.06), Inches(0.42), fill=col)
        mktx(s, title, fx + Inches(0.18), fy, Inches(1.4), Inches(0.28),
              size=12, bold=True, color=NAVY)
        mktx(s, desc, fx + Inches(0.18), fy + Inches(0.3), fw - Inches(0.2), Inches(0.38),
              size=10.5, color=BODY, wrap=True)
        mkrect(s, fx, fy + Inches(0.66), fw, Inches(0.012), fill=TEAL_XL)
        fy += Inches(0.68)

    # Portal mockup
    mkrect(s, Inches(6.25), Inches(1.55), Inches(6.7), Inches(5.2),
           fill=WHITE, line=NAVY, lw=Pt(1.5))
    mkrect(s, Inches(6.25), Inches(1.55), Inches(6.7), Inches(0.52), fill=NAVY)
    mktx(s, "  Sequor  Portal", Inches(6.25), Inches(1.62), Inches(4.5), Inches(0.38),
          size=11, bold=True, color=WHITE)

    stats = [("Messages", "1,284", "+12%", TEAL), ("Auto-Replied", "1,027", "80%", GREEN),
             ("Escalated", "42", "-3%", AMBER), ("Avg Time", "3.2 hr", "4hr SLA", NAVY)]
    sx = Inches(6.45)
    sy = Inches(2.2)
    for label, val, chg, col in stats:
        mkrect(s, sx, sy, Inches(1.42), Inches(0.78), fill=BG, line=TEAL_XL, lw=Pt(0.75))
        mkrect(s, sx, sy, Inches(1.42), Inches(0.055), fill=col)
        mktx(s, label, sx + Inches(0.1), sy + Inches(0.1), Inches(1.25), Inches(0.22),
              size=8.5, color=MUTED)
        mktx(s, val, sx + Inches(0.1), sy + Inches(0.3), Inches(1.0), Inches(0.32),
              size=16, bold=True, color=NAVY)
        sx += Inches(1.55)

    # Bar chart mock
    mkrect(s, Inches(6.45), Inches(3.12), Inches(3.65), Inches(1.55),
           fill=BG, line=TEAL_XL, lw=Pt(0.75))
    mktx(s, "Messages — Last 7 Days", Inches(6.57), Inches(3.22),
          Inches(2.5), Inches(0.25), size=9, bold=True, color=NAVY)
    bars = [(0.45, TEAL_L), (0.65, TEAL_L), (0.5, TEAL_L),
            (0.85, TEAL), (0.6, TEAL_L), (0.95, TEAL), (0.7, TEAL_L)]
    days = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
    bx = Inches(6.65)
    bw = Inches(0.33)
    for frac, col in bars:
        bh = frac * Inches(0.78)
        mkrect(s, bx, Inches(4.84) - bh, bw, bh, fill=col)
        bx += bw + Inches(0.06)

    # Recent messages mock
    mkrect(s, Inches(10.25), Inches(3.12), Inches(2.55), Inches(1.55),
           fill=BG, line=TEAL_XL, lw=Pt(0.75))
    mktx(s, "Recent", Inches(10.37), Inches(3.22), Inches(2), Inches(0.25),
          size=9, bold=True, color=NAVY)
    msgs = [("alice@co.com", "Order #12345", "Auto", GREEN),
            ("+1 555 0000", "Refund request", "Esc", AMBER),
            ("bob@store.io", "Shipping status", "Auto", GREEN)]
    my2 = Inches(3.5)
    for sender, txt, status, col in msgs:
        mktx(s, sender, Inches(10.37), my2, Inches(2.3), Inches(0.22),
              size=8.5, color=NAVY)
        mktx(s, txt, Inches(10.37), my2 + Inches(0.2), Inches(1.65), Inches(0.2),
              size=8, color=MUTED)
        pill(s, Inches(11.75), my2 + Inches(0.2), status, bg_col=col, w=Inches(0.7))
        my2 += Inches(0.44)


# ── SLIDE 5: AI Pipeline ───────────────────────────────────────────────────────
def s5():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, "How It Works: The AI Pipeline",
           "From inbound message to intelligent reply — in seconds, fully automated")
    footer_bar(s)

    steps = [
        (TEAL,   "1", "Message Received",
         "WhatsApp or Email arrives via webhook.\nMetadata and body captured and\ntimestamped, routed to pipeline."),
        (TEAL_L, "2", "Intent Classified",
         "LLM classifier identifies intent:\nrefund, order status, complaint,\nor general inquiry."),
        (NAVY,   "3", "Knowledge Retrieved",
         "RAG searches the Document Hub\nfor relevant content — pricing,\npolicies, guides, past cases."),
        (TEAL_L, "4", "Decision Made",
         "AutoReplyService evaluates:\nauto-reply or escalate?\nConfidence score sets routing."),
        (TEAL,   "5", "Reply Sent",
         "AI reply sent via SendGrid or\nWhatsApp API. Full audit log\nwritten to the database."),
    ]
    sw = Inches(2.2)
    sx = Inches(0.5)
    sy = Inches(1.85)
    for i, (col, num, title, desc) in enumerate(steps):
        mkrect(s, sx, sy, sw, Inches(3.55), fill=col)
        mkrect(s, sx + Inches(0.8), sy + Inches(0.15),
               Inches(0.6), Inches(0.6), fill=WHITE)
        mktx(s, num, sx + Inches(0.8), sy + Inches(0.17),
              Inches(0.6), Inches(0.6),
              size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
        mktx(s, title, sx, sy + Inches(0.88), sw, Inches(0.5),
              size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        mktx(s, desc, sx + Inches(0.1), sy + Inches(1.48),
              sw - Inches(0.2), Inches(1.9),
              size=11.5, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)
        if i < len(steps) - 1:
            mktx(s, "→", sx + sw + Inches(0.03), sy + Inches(1.3),
                  Inches(0.42), Inches(0.5),
                  size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        sx += sw + Inches(0.52)

    mkrect(s, Inches(0.5), Inches(5.65), Inches(12.3), Inches(1.22),
           fill=WHITE, line=TEAL_XL, lw=Pt(0.75))
    mkrect(s, Inches(0.5), Inches(5.65), Inches(0.07), Inches(1.22), fill=GREEN)
    mktx(s, "Continuous Learning Loop",
          Inches(0.75), Inches(5.77), Inches(4.5), Inches(0.35),
          size=13, bold=True, color=NAVY)
    mktx(s, "After every escalation is resolved, the LearningLoop records the operator's feedback — "
             "improving future classification accuracy and reply quality. Every resolved case makes "
             "the AI smarter for all customers.",
          Inches(0.75), Inches(6.15), Inches(11.7), Inches(0.6),
          size=11.5, color=BODY, wrap=True)


# ── SLIDE 6: Channels ───────────────────────────────────────────────────────
def s6():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, "Multi-Channel: WhatsApp + Email",
           "Both channels share the same AI pipeline, knowledge base, and operator portal")
    footer_bar(s)

    # WhatsApp
    mkrect(s, Inches(0.62), Inches(1.68), Inches(5.75), Inches(4.65),
           fill=WHITE, line=TEAL_XL, lw=Pt(0.75))
    mkrect(s, Inches(0.62), Inches(1.68), Inches(5.75), Inches(0.65), fill=NAVY)
    mktx(s, "WhatsApp", Inches(0.88), Inches(1.82), Inches(4.5), Inches(0.42),
          size=17, bold=True, color=WHITE)
    wa = [
        "Meta Cloud API — modern, stable, officially supported",
        "Webhook verification with HMAC-SHA256 signatures",
        "Template messages for outbound notifications",
        "Contact profiling and phone number validation",
        "Media: images, documents, audio, voice notes",
        "Delivery and read receipts tracked in portal",
        "Opt-out and block handling fully automated",
        "Business profile management via API",
    ]
    wy = Inches(2.52)
    for f in wa:
        mktx(s, f"  {f}", Inches(0.88), wy, Inches(5.3), Inches(0.34),
              size=12, color=BODY)
        wy += Inches(0.46)

    # Email
    mkrect(s, Inches(6.65), Inches(1.68), Inches(6.1), Inches(4.65),
           fill=WHITE, line=TEAL_XL, lw=Pt(0.75))
    mkrect(s, Inches(6.65), Inches(1.68), Inches(6.1), Inches(0.65), fill=TEAL)
    mktx(s, "Email", Inches(6.9), Inches(1.82), Inches(4.5), Inches(0.42),
          size=17, bold=True, color=WHITE)
    em = [
        "SendGrid Inbound Parse — full email receipt with attachments",
        "Attachment ingestion: files saved to Document Hub automatically",
        "HTML and plain-text body parsing, intelligently normalized",
        "Reply-to threading via In-Reply-To and References headers",
        "SPF / DKIM / DMARC aware routing and validation",
        "Auto-reply with quoted original for customer clarity",
        "PDPA: email content minimized in storage post-resolution",
        "Bounce handling and suppression list management",
    ]
    ey = Inches(2.52)
    for f in em:
        mktx(s, f"  {f}", Inches(6.9), ey, Inches(5.65), Inches(0.34),
              size=12, color=BODY)
        ey += Inches(0.46)

    mktx(s, "Both channels: same AI pipeline — classify → RAG → decide → reply",
          Inches(0.62), Inches(6.48), Inches(12.1), Inches(0.35),
          size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ── SLIDE 7: Escalations ─────────────────────────────────────────────────────
def s7():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, "Smart Escalations",
           "The right message, to the right operator, with full context — instantly")
    footer_bar(s)

    mktx(s, "Escalation Triggers", Inches(0.62), Inches(1.68), Inches(5.5), Inches(0.38),
          size=14, bold=True, color=NAVY)

    triggers = [
        (RED,   "Cancellation / Refund Threats",
         "High-friction intent signals — classifier confidence > 0.8 triggers immediate escalation."),
        (AMBER, "VIP / Priority Contacts",
         "Contacts flagged by tier or manually tagged receive instant priority routing."),
        (AMBER, "SLA Breach Risk",
         "Messages unresolved past the configured SLA window auto-escalate."),
        (RED,   "Negative Sentiment",
         "Frustration or anger detected by the classifier — regardless of content topic."),
        (RED,   "Repeat Queries (3x+)",
         "Same topic from same contact, unanswered, escalates after 3 attempts."),
        (TEAL,  "Manual Flag",
         "Operator can flag any thread for escalation with one click and a note."),
    ]
    ty = Inches(2.1)
    for col, title, desc in triggers:
        mkrect(s, Inches(0.62), ty + Inches(0.06), Inches(0.06), Inches(0.42), fill=col)
        mktx(s, title, Inches(0.82), ty, Inches(2.5), Inches(0.28),
              size=12, bold=True, color=NAVY)
        mktx(s, desc, Inches(0.82), ty + Inches(0.3), Inches(5.2), Inches(0.4),
              size=10.5, color=BODY, wrap=True)
        ty += Inches(0.73)

    mkrect(s, Inches(6.65), Inches(1.58), Inches(6.15), Inches(5.18),
           fill=WHITE, line=TEAL_XL, lw=Pt(0.75))
    mkrect(s, Inches(6.65), Inches(1.58), Inches(6.15), Inches(0.56), fill=NAVY)
    mktx(s, "Resolution Flow", Inches(6.9), Inches(1.72), Inches(5.5), Inches(0.36),
          size=14, bold=True, color=WHITE)

    flow = [
        ("1", "Alert Created",
         "Escalation record created with full message context, contact info, and SLA countdown."),
        ("2", "Assigned",
         "Auto-assigned to available operator based on team routing rules and contact tier."),
        ("3", "Context Delivered",
         "Portal shows full thread plus relevant document excerpts from RAG."),
        ("4", "Operator Resolves",
         "Outcome logged. LearningLoop records resolution for training data."),
        ("5", "SLA Closed",
         "Ticket closed. Resolution time and outcome stored in immutable audit trail."),
    ]
    fy = Inches(2.3)
    for num, title, desc in flow:
        mkrect(s, Inches(6.88), fy, Inches(0.5), Inches(0.5), fill=TEAL)
        mktx(s, num, Inches(6.88), fy + Inches(0.04), Inches(0.5), Inches(0.45),
              size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        mktx(s, title, Inches(7.55), fy + Inches(0.04), Inches(2.8), Inches(0.3),
              size=12, bold=True, color=NAVY)
        mktx(s, desc, Inches(7.55), fy + Inches(0.35), Inches(5.0), Inches(0.45),
              size=10.5, color=BODY, wrap=True)
        if num != "5":
            mkrect(s, Inches(7.1), fy + Inches(0.55), Inches(0.06), Inches(0.2), fill=TEAL_XL)
        fy += Inches(0.85)


# ── SLIDE 8: PDPA ─────────────────────────────────────────────────────────────
def s8():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, "PDPA Compliance by Design",
           "Every layer of Sequor is built with Singapore data protection requirements at its core")
    footer_bar(s)

    pillars = [
        (TEAL,   "1", "Data Minimization",
         "Sequor stores only operationally necessary data. Email bodies and WhatsApp content are "
         "retained only for active threads. Historical archives are pseudonymized."),
        (NAVY,   "2", "Right to Erasure",
         "Full deletion workflow: data subject requests erasure, all PII purged from active "
         "records and audit logs, erasure confirmation issued automatically."),
        (TEAL_L, "3", "Immutable Audit Trail",
         "Every message, classification, escalation, and operator action logged with tenant_id, "
         "timestamp, actor, and outcome — PDPA Article 20 compliant."),
        (TEAL,   "4", "Tenant Isolation",
         "Role-based access: operators see only assigned contacts. Database-level tenant isolation "
         "enforced on every query. No cross-tenant data leakage."),
        (NAVY,   "5", "Consent Management",
         "All inbound contacts flagged for consent status. Audit log records channel, notice "
         "version, and timestamp — retrievable on demand."),
        (TEAL_L, "6", "PII Classification",
         "Email addresses, phone numbers, and names classified as PII in the audit trail. RAG "
         "retrieval respects classification — no PII in AI training data."),
    ]
    pw = Inches(3.85)
    px = Inches(0.62)
    py = Inches(1.72)
    for i, (col, num, title, desc) in enumerate(pillars):
        row = i // 3
        col_i = i % 3
        x = px + col_i * (pw + Inches(0.28))
        y = py + row * Inches(2.6)
        mkrect(s, x, y, pw, Inches(2.4), fill=WHITE, line=TEAL_XL, lw=Pt(0.75))
        mkrect(s, x, y, pw, Inches(0.07), fill=col)
        mkrect(s, x + Inches(0.2), y + Inches(0.18), Inches(0.5), Inches(0.5), fill=col)
        mktx(s, num, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5),
              size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        mktx(s, title, x + Inches(0.85), y + Inches(0.25), pw - Inches(1.0), Inches(0.38),
              size=13, bold=True, color=NAVY)
        mktx(s, desc, x + Inches(0.2), y + Inches(0.82), pw - Inches(0.4), Inches(1.45),
              size=11, color=BODY, wrap=True)


# ── SLIDE 9: Pricing ───────────────────────────────────────────────────────────
def s9():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, "Pricing",
           "Simple, transparent, scalable — free forever for development and testing")
    footer_bar(s)

    plans = [
        ("Free",    "$0",   "month",
         ["50 messages / month", "WhatsApp + Email channels",
          "Auto-reply AI", "RAG Document Hub",
          "Dashboard & message history",
          "1 operator account",
          "Community support"],
         TEAL_XL, NAVY, False),
        ("Solo",    "$15",  "month",
         ["200 messages / month", "WhatsApp + Email channels",
          "Auto-reply AI + RAG",
          "Dashboard & message history",
          "30-day message retention",
          "1 operator account",
          "Email support"],
         TEAL_XL, NAVY, False),
        ("Starter", "$35",  "month",
         ["Unlimited messages*", "WhatsApp + Email channels",
          "Auto-reply AI + RAG",
          "Smart Escalations",
          "PDPA audit trail",
          "90-day message retention",
          "Up to 3 operators",
          "Priority support"],
         TEAL, WHITE, True),
        ("Pro",     "$55",  "month",
         ["Unlimited messages",
          "Everything in Starter",
          "Unlimited operators",
          "Unlimited documents",
          "Priority RAG processing",
          "Advanced analytics dashboard",
          "Custom key phrase maps",
          "PDPA compliance report",
          "Dedicated support"],
         NAVY, WHITE, False),
    ]
    pw = Inches(2.9)
    px = Inches(0.55)
    py = Inches(1.8)
    for name, price, period, features, col, txt, featured in plans:
        ph = Inches(5.1)
        mkrect(s, px, py, pw, ph, fill=col,
               line=col if featured else TEAL_XL,
               lw=Pt(2) if featured else Pt(1))
        if featured:
            mkrect(s, px, py, pw, Inches(0.42), fill=TEAL_L)
            mktx(s, "Most Popular", px, py + Inches(0.08), pw, Inches(0.32),
                  size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            ty = py + Inches(0.55)
        else:
            ty = py + Inches(0.22)
        mktx(s, name, px, ty, pw, Inches(0.55),
              size=22, bold=True, color=txt, align=PP_ALIGN.CENTER)
        mktx(s, price, px, ty + Inches(0.6), pw, Inches(0.75),
              size=40, bold=True, color=txt, align=PP_ALIGN.CENTER)
        mktx(s, f"/{period}", px, ty + Inches(1.32), pw, Inches(0.3),
              size=12, color=txt, align=PP_ALIGN.CENTER)
        hline(s, px + Inches(0.28), ty + Inches(1.72),
              pw - Inches(0.56), color=txt if featured else TEAL_XL, thick=Pt(0.6))
        fy = ty + Inches(1.85)
        for feat in features:
            mktx(s, f"  {feat}", px + Inches(0.22), fy,
                  pw - Inches(0.38), Inches(0.36),
                  size=11, color=txt)
            fy += Inches(0.4)
        px += pw + Inches(0.22)

    mktx(s, "* Message fair-use limits enforced at 3x plan average. Enterprise pricing available on request.",
          Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.3),
          size=10, color=MUTED, italic=True)


# ── SLIDE 10: Setup ────────────────────────────────────────────────────────────
def s10():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, "Getting Started",
           "No servers, no ML pipelines, no data science team required")
    footer_bar(s)

    steps = [
        (TEAL,   "1", "Connect WhatsApp Business Account",
         "Log in to Meta Business Manager → add the WhatsApp Business API. "
         "Copy your phone number ID and access token into the Sequor portal. Takes 3 minutes."),
        (NAVY,   "2", "Connect Email via SendGrid",
         "Create a free SendGrid account → configure Inbound Parse webhook → "
         "point it at your Sequor URL. Copy your API key in. Done."),
        (TEAL_L, "3", "Upload Your Documents",
         "Drop your PDFs, policy docs, FAQs, and product guides into the Document Hub. "
         "RAG indexing starts automatically — approximately 5 minutes for 50 documents."),
        (TEAL,   "4", "Configure Key Phrases & Escalation Rules",
         "Add topics your customers ask about. Map them to reply strategies. "
         "Set your SLA window and escalation contacts. Roughly 10 minutes."),
        (NAVY,   "5", "Go Live",
         "Flip the channel to Active. Every inbound message now routes through the AI pipeline. "
         "First auto-reply typically fires within the hour."),
    ]
    sw = Inches(5.8)
    sx = Inches(0.62)
    sy = Inches(1.78)
    for col, num, title, desc in steps:
        mkrect(s, sx, sy, Inches(0.55), Inches(0.9), fill=col)
        mktx(s, num, sx, sy + Inches(0.04), Inches(0.55), Inches(0.9),
              size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        mktx(s, title, sx + Inches(0.68), sy, sw - Inches(0.68), Inches(0.35),
              size=13, bold=True, color=NAVY)
        mktx(s, desc, sx + Inches(0.68), sy + Inches(0.38), sw - Inches(0.68), Inches(0.52),
              size=11, color=BODY, wrap=True)
        sy += Inches(1.04)

    mkrect(s, Inches(6.9), Inches(1.68), Inches(5.95), Inches(5.1),
           fill=WHITE, line=TEAL_XL, lw=Pt(0.75))
    mkrect(s, Inches(6.9), Inches(1.68), Inches(5.95), Inches(0.6), fill=NAVY)
    mktx(s, "About Document Quality", Inches(7.15), Inches(1.82),
          Inches(5.5), Inches(0.4), size=15, bold=True, color=WHITE)
    mktx(s,
          "The AI's accuracy is directly proportional to your document quality. "
          "Well-structured PDFs — Q&A format, clear headings, policy docs — "
          "produce the best auto-replies. Poorly formatted scans perform accordingly.",
          Inches(7.15), Inches(2.45), Inches(5.5), Inches(1.05),
          size=12, color=BODY, wrap=True)
    hline(s, Inches(7.15), Inches(3.6), Inches(5.5))
    mktx(s, "Best document types:", Inches(7.15), Inches(3.72),
          Inches(5.5), Inches(0.3), size=12, bold=True, color=NAVY)
    best = [
        "Product FAQ sheets (Q&A format)",
        "Policy documents (pricing, shipping, returns)",
        "Service descriptions and menus",
        "How-to guides and troubleshooting docs",
        "Pricing sheets and plan comparisons",
    ]
    by = Inches(4.1)
    for b in best:
        mktx(s, f"  {b}", Inches(7.15), by, Inches(5.5), Inches(0.32),
              size=11.5, color=BODY)
        by += Inches(0.38)


# ── SLIDE 11: Demo ───────────────────────────────────────────────────────────
def s11():
    s = prs.slides.add_slide(BLANK)
    mkrect(s, 0, 0, W, H, fill=DARK)
    mkrect(s, Inches(8.2), 0, Inches(5.13), H, fill=NAVY)
    mkrect(s, Inches(9.7), 0, Inches(3.63), H, fill=NAVY2)

    mktx(s, "Live Demo", Inches(0.75), Inches(0.95), Inches(7), Inches(0.85),
          size=42, bold=True, color=WHITE)
    hline(s, Inches(0.75), Inches(1.88), Inches(3.5))
    mktx(s, "Watch Sequor handle a real inbound message end-to-end",
          Inches(0.75), Inches(2.0), Inches(7.5), Inches(0.42),
          size=15, color=TEAL_L)

    demo = [
        ("1", "Inbound message arrives", "via WhatsApp or Email webhook"),
        ("2", "AI classifies intent", "and retrieves relevant knowledge from Document Hub"),
        ("3", "Auto-reply generated", "and sent — or escalation triggered to operator"),
        ("4", "Portal updates", "dashboard, message history, and escalation queue"),
        ("5", "Escalation resolved", "operator closes ticket, AI learns for next time"),
    ]
    dx = Inches(0.75)
    dy = Inches(2.95)
    for num, step, detail in demo:
        mkrect(s, dx, dy, Inches(0.55), Inches(0.55), fill=TEAL)
        mktx(s, num, dx, dy + Inches(0.04), Inches(0.55), Inches(0.5),
              size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        mktx(s, step, dx + Inches(0.72), dy + Inches(0.04), Inches(3.2), Inches(0.32),
              size=13, bold=True, color=WHITE)
        mktx(s, detail, dx + Inches(0.72), dy + Inches(0.35), Inches(6.5), Inches(0.28),
              size=12, color=TEAL_L)
        dy += Inches(0.75)

    mktx(s, "Try it now  →  portal.sequor.com",
          Inches(0.75), Inches(6.6), Inches(7), Inches(0.38),
          size=13, bold=True, color=TEAL)

    mktx(s, "Portal\nLive",
          Inches(9.0), Inches(2.3), Inches(3.5), Inches(1.4),
          size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    mktx(s, "Sequor Pte. Ltd.\nSingapore",
          Inches(9.0), Inches(3.85), Inches(3.5), Inches(0.7),
          size=12, color=TEAL_L, align=PP_ALIGN.CENTER)


# ── SLIDE 12: Roadmap ────────────────────────────────────────────────────────
def s12():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, "Roadmap",
           "What's coming — building the full customer communication platform")
    footer_bar(s)

    quarters = [
        ("Q3 2026", TEAL, [
            ("SMS Channel",
             "Twilio integration for SMS inbound/outbound — complete omnichannel coverage."),
            ("Team Roles & Permissions",
             "Admin, Manager, Operator tiers with granular access controls."),
            ("Analytics Dashboard",
             "Message volume trends, AI accuracy reports, SLA performance tracking."),
        ]),
        ("Q4 2026", NAVY, [
            ("Multi-Language AI",
             "Thai, Bahasa, Vietnamese, Mandarin — AI replies in the customer's language."),
            ("API Webhooks",
             "Custom outbound triggers for CRM and ERP integration."),
            ("Mobile App",
             "iOS and Android — escalation management and approvals on the go."),
        ]),
        ("Q1 2027", NAVY2, [
            ("Enterprise SSO",
             "SAML and OIDC single sign-on for large enterprise procurement."),
            ("SOC 2 Type II",
             "Security certification for Fortune 500 and regulated industry procurement."),
            ("Voice Channel",
             "Inbound voice with transcription and AI summarization before routing."),
        ]),
    ]
    rw = Inches(3.9)
    rx = Inches(0.62)
    ry = Inches(1.8)
    for qtr, col, items in quarters:
        mkrect(s, rx, ry, rw, Inches(5.1), fill=WHITE, line=TEAL_XL, lw=Pt(0.75))
        mkrect(s, rx, ry, rw, Inches(0.65), fill=col)
        mktx(s, qtr, rx, ry + Inches(0.14), rw, Inches(0.45),
              size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        iy = ry + Inches(0.88)
        for title, desc in items:
            mkrect(s, rx + Inches(0.22), iy + Inches(0.08),
                   Inches(0.06), Inches(0.42), fill=col)
            mktx(s, title, rx + Inches(0.42), iy, rw - Inches(0.6), Inches(0.35),
                  size=13, bold=True, color=NAVY)
            mktx(s, desc, rx + Inches(0.42), iy + Inches(0.38),
                  rw - Inches(0.6), Inches(0.6), size=11.5, color=BODY, wrap=True)
            iy += Inches(1.22)
        rx += rw + Inches(0.28)


# ── SLIDE 13: Why Sequor ──────────────────────────────────────────────────────
def s13():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, "Why Sequor",
           "Built for Southeast Asian SMBs who need enterprise-grade AI without the complexity")
    footer_bar(s)

    reasons = [
        (TEAL,   "Designed for Southeast Asia",
         "WhatsApp as the primary channel, multi-language AI, PDPA-first compliance — "
         "designed for the region's SMB reality, not adapted from Western enterprise software."),
        (NAVY,   "10x Operator Efficiency",
         "80% of messages auto-replied. Operators spend time on the 20% that matters — "
         "escalations, VIPs, and complex cases — not repetitive FAQs."),
        (TEAL_L, "Zero Infrastructure Headaches",
         "No servers, no ML pipelines, no training data to build. Connect your channels "
         "and upload your documents. The AI handles the rest."),
        (TEAL,   "PDPA from Day One",
         "Compliance isn't an afterthought or an add-on module. Audit trails, erasure "
         "workflows, and data minimization are built into every layer."),
        (NAVY,   "Transparent Pricing",
         "Predictable per-seat pricing. No per-message AI credits, no hidden ingestion "
         "fees, no enterprise-only features locked behind a sales call."),
        (TEAL_L, "Fast Time to Value",
         "Channels connected in minutes. First AI reply fires within the hour. "
         "Document Hub built from your existing files — no structured data required."),
    ]
    rw = Inches(3.9)
    rx = Inches(0.62)
    ry = Inches(1.8)
    for i, (col, title, desc) in enumerate(reasons):
        row = i // 3
        col_i = i % 3
        x = rx + col_i * (rw + Inches(0.28))
        y = ry + row * Inches(2.55)
        mkrect(s, x, y, rw, Inches(2.3), fill=WHITE, line=TEAL_XL, lw=Pt(0.75))
        mkrect(s, x, y, rw, Inches(0.07), fill=col)
        mktx(s, title, x + Inches(0.22), y + Inches(0.2), rw - Inches(0.4), Inches(0.5),
              size=13.5, bold=True, color=NAVY)
        mktx(s, desc, x + Inches(0.22), y + Inches(0.78), rw - Inches(0.4), Inches(1.4),
              size=12, color=BODY, wrap=True)


# ── SLIDE 14: Opportunity ────────────────────────────────────────────────────
def s14():
    s = prs.slides.add_slide(BLANK)
    mkrect(s, 0, 0, W, H, fill=DARK)
    mkrect(s, Inches(8.2), 0, Inches(5.13), H, fill=NAVY)
    mkrect(s, Inches(9.7), 0, Inches(3.63), H, fill=NAVY2)

    mktx(s, "The Opportunity", Inches(0.75), Inches(0.65), Inches(7.5), Inches(0.85),
          size=36, bold=True, color=WHITE)
    hline(s, Inches(0.75), Inches(1.6), Inches(3.5))

    stats = [
        ("$2.4B", "TAM — SE Asian SMB Customer\nCommunication Software (2026)"),
        ("35%",   "Annual growth in WhatsApp\nBusiness API adoption"),
        ("80%",   "Average auto-reply rate\nin Sequor pilot accounts"),
        ("< 4 wks", "Customer time-to-first\nAI-replied message"),
    ]
    sx = Inches(0.75)
    sy = Inches(2.0)
    for val, desc in stats:
        mktx(s, val, sx, sy, Inches(2.6), Inches(0.82),
              size=34, bold=True, color=TEAL)
        mktx(s, desc, sx, sy + Inches(0.88), Inches(2.6), Inches(0.75),
              size=11.5, color=TEAL_L, wrap=True)
        sy += Inches(1.85)

    mkrect(s, Inches(8.2), Inches(1.58), Inches(4.5), Inches(5.18), fill=NAVY)
    mkrect(s, Inches(8.2), Inches(1.58), Inches(4.5), Inches(0.65), fill=TEAL)
    mktx(s, "Investment Summary",
          Inches(8.4), Inches(1.72), Inches(4.1), Inches(0.42),
          size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("Raise", "Seed Round — S$1.2M"),
        ("Use of Funds", "Engineering (40%)\nSales & Marketing (35%)\nOperations (25%)"),
        ("Key Milestones", "20 paying customers\n1,000 monthly messages\n3 channel integrations"),
        ("Target Close", "Q3 2026"),
    ]
    ry2 = Inches(2.5)
    for label, val in rows:
        mktx(s, label + ":", Inches(8.4), ry2, Inches(1.5), Inches(0.3),
              size=11, bold=True, color=TEAL_L)
        mktx(s, val, Inches(9.95), ry2, Inches(2.6), Inches(0.65),
              size=11.5, color=WHITE, wrap=True)
        ry2 += Inches(0.84)

    footer_bar(s)


# ── SLIDE 15: Contact ─────────────────────────────────────────────────────────
def s15():
    s = prs.slides.add_slide(BLANK)
    mkrect(s, 0, 0, W, H, fill=DARK)
    mkrect(s, Inches(8.2), 0, Inches(5.13), H, fill=NAVY)
    mkrect(s, Inches(9.7), 0, Inches(3.63), H, fill=NAVY2)

    mkrect(s, Inches(0.75), Inches(1.4), Inches(0.85), Inches(0.85), fill=TEAL)
    mktx(s, "S", Inches(0.75), Inches(1.42), Inches(0.85), Inches(0.85),
          size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    mktx(s, "Sequor", Inches(1.78), Inches(1.47), Inches(4), Inches(0.75),
          size=38, bold=True, color=WHITE)
    hline(s, Inches(0.75), Inches(2.4), Inches(5))
    mktx(s, "AI-Powered Customer Communication Platform",
          Inches(0.75), Inches(2.52), Inches(7), Inches(0.38),
          size=14, color=TEAL_L)

    contacts = [
        ("Website:",  "www.sequor.com"),
        ("Email:",   "hello@sequor.com"),
        ("WhatsApp:", "+65 9000 0000"),
        ("Address:", "1 Raffles Place, #12-00\nSingapore 048616"),
    ]
    cy = Inches(3.2)
    for label, val in contacts:
        mktx(s, label, Inches(0.75), cy, Inches(2.2), Inches(0.32),
              size=13, bold=True, color=TEAL_L)
        mktx(s, val, Inches(2.95), cy, Inches(5), Inches(0.45),
              size=13, color=WHITE)
        cy += Inches(0.6)

    mktx(s, "Thank you for\nyour time.",
          Inches(9.0), Inches(2.3), Inches(3.5), Inches(1.4),
          size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    hline(s, Inches(9.0), Inches(3.85), Inches(3.5))
    mktx(s, "We look forward to building\nthe future of customer communication.",
          Inches(9.0), Inches(4.05), Inches(3.5), Inches(0.85),
          size=12, color=TEAL_L, align=PP_ALIGN.CENTER, wrap=True)

    footer_bar(s)


# ── Build ─────────────────────────────────────────────────────────────────────
s1();  s2();  s3();  s4();  s5()
s6();  s7();  s8();  s9();  s10()
s11(); s12(); s13(); s14(); s15()

out = "/Users/aliciapang/Documents/GitHub/Sequor/workspaces/_template/Sequor_Investor_Deck_May2026.pptx"
prs.save(out)
print(f"Saved: {out}")
