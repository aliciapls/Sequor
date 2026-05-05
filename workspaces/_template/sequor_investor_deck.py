"""Generate Sequor PE Investor Presentation Deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette (duck egg blue) ────────────────────────────────────────────────────
D50  = RGBColor(0xF0, 0xF9, 0xFB)   # slide backgrounds
D100 = RGBColor(0xD9, 0xEE, 0xF4)
D200 = RGBColor(0xB7, 0xDD, 0xEA)
D300 = RGBColor(0x8F, 0xC5, 0xDB)
D400 = RGBColor(0x5F, 0xA6, 0xC5)
D500 = RGBColor(0x3C, 0x8E, 0xAF)   # primary accent
D600 = RGBColor(0x2D, 0x72, 0x96)
D700 = RGBColor(0x1D, 0x5A, 0x77)   # dark text / headers
D800 = RGBColor(0x1A, 0x4A, 0x63)   # darkest
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED   = RGBColor(0xEF, 0x44, 0x44)

# ── Slide dimensions (widescreen 16:9) ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helper utilities ───────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_width
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width=None, height=None,
             font_size=Pt(18), bold=False, color=D700,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    from pptx.util import Pt
    if width and height:
        txb = slide.shapes.add_textbox(left, top, width, height)
    else:
        txb = slide.shapes.add_textbox(left, top, Inches(1), Inches(1))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Inter"
    return txb


def add_para(tf, text, font_size=Pt(16), bold=False, color=D700,
             align=PP_ALIGN.LEFT, italic=False, space_before=Pt(0)):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Inter"
    return p


def card(slide, left, top, width, height, title, body_lines,
         icon=None, accent=D500):
    """Feature card with icon area, title, and bullet lines."""
    # Card background
    add_rect(slide, left, top, width, height, fill_rgb=WHITE,
             line_rgb=D200, line_width=Pt(1))
    # Accent top bar
    add_rect(slide, left, top, width, Inches(0.06), fill_rgb=accent)
    # Title
    add_text(slide, title, left + Inches(0.25), top + Inches(0.2),
             width=width - Inches(0.5), font_size=Pt(15), bold=True, color=D800)
    # Body lines
    y = top + Inches(0.6)
    for line in body_lines:
        add_text(slide, f"• {line}", left + Inches(0.25), y,
                 width=width - Inches(0.5), font_size=Pt(12), color=SLATE)
        y += Inches(0.35)


# ── SLIDE 1: Cover ─────────────────────────────────────────────────────────────
def slide_cover():
    slide = prs.slides.add_slide(BLANK)
    # Full dark gradient background
    add_rect(slide, 0, 0, W, H, fill_rgb=D800)
    # Decorative circle top-right
    add_rect(slide, Inches(9.5), Inches(-1.5), Inches(5), Inches(5), fill_rgb=D700)
    add_rect(slide, Inches(10.5), Inches(-0.5), Inches(4), Inches(4), fill_rgb=D600)
    # Decorative circle bottom-left
    add_rect(slide, Inches(-1.5), Inches(5), Inches(4), Inches(4), fill_rgb=D700)
    add_rect(slide, Inches(-0.5), Inches(5.5), Inches(3), Inches(3), fill_rgb=D600)
    # Logo box
    add_rect(slide, Inches(1), Inches(1.8), Inches(0.9), Inches(0.9),
             fill_rgb=D500)
    add_text(slide, "S", Inches(1), Inches(1.82), Inches(0.9), Inches(0.9),
             font_size=Pt(36), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Brand name
    add_text(slide, "Sequor", Inches(2.1), Inches(1.85), Inches(3), Inches(0.8),
             font_size=Pt(40), bold=True, color=WHITE)
    # Tagline
    add_text(slide, "AI-Powered Customer Communication Platform",
             Inches(1), Inches(2.85), Inches(9), Inches(0.6),
             font_size=Pt(22), bold=False, color=D300)
    # Divider line
    add_rect(slide, Inches(1), Inches(3.55), Inches(6), Inches(0.04), fill_rgb=D500)
    # Subtitle
    add_text(slide, "Investor Presentation", Inches(1), Inches(3.75), Inches(6),
             font_size=Pt(16), bold=False, color=D400)
    # Bottom bar
    add_rect(slide, 0, Inches(6.7), W, Inches(0.8), fill_rgb=D700)
    add_text(slide, "Confidential  •  May 2026", Inches(1), Inches(6.78),
             Inches(5), Inches(0.5), font_size=Pt(11), color=D400, italic=True)


# ── SLIDE 2: Problem ────────────────────────────────────────────────────────────
def slide_problem():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=D50)
    # Header bar
    add_rect(slide, 0, 0, W, Inches(1.3), fill_rgb=D800)
    add_text(slide, "The Problem", Inches(0.6), Inches(0.3), Inches(6),
             font_size=Pt(30), bold=True, color=WHITE)
    add_text(slide, "Why businesses are losing customers every day",
             Inches(0.6), Inches(0.82), Inches(8), Inches(0.4),
             font_size=Pt(14), color=D300, italic=True)
    # Problem cards
    problems = [
        ("💬", "Channels Explosion",
         "Customers message on WhatsApp, email, SMS — but teams manage them all in silos, losing context at every handoff."),
        ("⏰", "Response Time Kills SLAs",
         "Average first response time: 4–6 hours. 35% of customers abandon the brand after one slow response."),
        ("🧠", "Repetitive Queries Exhaust Teams",
         "70% of inbound messages are repeats — product info, order status, FAQs. Agents burn out answering the same questions."),
        ("📋", "Escalations Fall Through Cracks",
         "High-intent signals (frustration, cancellation threats, VIPs) are missed. Escalations require 6+ tools to resolve."),
        ("🔒", "Compliance Burden",
         "PDPA requires audit trails, data minimization, and erasure workflows. Manual compliance is error-prone and costly."),
        ("📈", "Zero Visibility",
         "No unified view of message volume, resolution rates, or AI performance. Data lives in disconnected inboxes."),
    ]
    cols = 3
    card_w = Inches(3.8)
    card_h = Inches(1.6)
    start_x = Inches(0.5)
    start_y = Inches(1.6)
    gap_x = Inches(0.3)
    gap_y = Inches(0.25)
    for i, (icon, title, desc) in enumerate(problems):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        # Card bg
        add_rect(slide, x, y, card_w, card_h, fill_rgb=WHITE, line_rgb=D200, line_width=Pt(1))
        # Icon + title row
        add_rect(slide, x, y, Inches(0.55), card_h, fill_rgb=D100)
        add_text(slide, icon, x + Inches(0.05), y + Inches(0.55),
                 Inches(0.45), Inches(0.5), font_size=Pt(22), align=PP_ALIGN.CENTER)
        add_text(slide, title, x + Inches(0.6), y + Inches(0.15),
                 card_w - Inches(0.75), font_size=Pt(13), bold=True, color=D800)
        add_text(slide, desc, x + Inches(0.6), y + Inches(0.5),
                 card_w - Inches(0.75), font_size=Pt(10.5), color=SLATE, wrap=True)


# ── SLIDE 3: Solution ──────────────────────────────────────────────────────────
def slide_solution():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=D50)
    add_rect(slide, 0, 0, W, Inches(1.3), fill_rgb=D800)
    add_text(slide, "Sequor: The Solution", Inches(0.6), Inches(0.3), Inches(8),
             font_size=Pt(30), bold=True, color=WHITE)
    add_text(slide, "One platform — every channel, every message, fully automated",
             Inches(0.6), Inches(0.82), Inches(10), Inches(0.4),
             font_size=Pt(14), color=D300, italic=True)
    # Central value prop
    add_rect(slide, Inches(0.6), Inches(1.6), Inches(5.5), Inches(2.8), fill_rgb=D800)
    add_text(slide, "Sequor unifies WhatsApp, Email, and future channels into a single AI-powered platform that auto-replies, escalates intelligently, and gives operators complete visibility — while staying PDPA compliant by design.",
             Inches(0.9), Inches(1.9), Inches(5), Inches(2.4),
             font_size=Pt(15), color=WHITE, wrap=True)
    add_text(slide, "80%", Inches(6.8), Inches(1.7), Inches(1.5), Inches(1),
             font_size=Pt(52), bold=True, color=D500, align=PP_ALIGN.CENTER)
    add_text(slide, "of messages auto-replied", Inches(6.5), Inches(2.7), Inches(2.1), Inches(0.5),
             font_size=Pt(11), color=SLATE, align=PP_ALIGN.CENTER)
    add_text(slide, "4-hr", Inches(8.8), Inches(1.7), Inches(1.5), Inches(1),
             font_size=Pt(52), bold=True, color=D500, align=PP_ALIGN.CENTER)
    add_text(slide, "avg SLA achievement", Inches(8.5), Inches(2.7), Inches(2.1), Inches(0.5),
             font_size=Pt(11), color=SLATE, align=PP_ALIGN.CENTER)
    add_text(slide, "2 min", Inches(10.8), Inches(1.7), Inches(1.5), Inches(1),
             font_size=Pt(52), bold=True, color=D500, align=PP_ALIGN.CENTER)
    add_text(slide, "to set up a channel", Inches(10.5), Inches(2.7), Inches(2.1), Inches(0.5),
             font_size=Pt(11), color=SLATE, align=PP_ALIGN.CENTER)
    # Three pillars
    pillars = [
        ("Auto-Reply AI", "Classifies intent, retrieves relevant documents, and generates accurate replies — without human input.", D500),
        ("Smart Escalations", "Detects high-intent signals and routes to the right operator with full context.", D600),
        ("PDPA Compliance", "Built-in audit trail, data minimization, and erasure workflows — compliant by default.", D700),
    ]
    pw = Inches(3.8)
    px = Inches(0.6)
    py = Inches(4.6)
    for title, desc, col in pillars:
        add_rect(slide, px, py, pw, Inches(1.7), fill_rgb=WHITE, line_rgb=D200, line_width=Pt(1))
        add_rect(slide, px, py, pw, Inches(0.06), fill_rgb=col)
        add_text(slide, title, px + Inches(0.25), py + Inches(0.2),
                 pw - Inches(0.5), font_size=Pt(14), bold=True, color=D800)
        add_text(slide, desc, px + Inches(0.25), py + Inches(0.55),
                 pw - Inches(0.5), font_size=Pt(11.5), color=SLATE, wrap=True)
        px += pw + Inches(0.25)


# ── SLIDE 4: Product Overview ──────────────────────────────────────────────────
def slide_product():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=D50)
    add_rect(slide, 0, 0, W, Inches(1.3), fill_rgb=D800)
    add_text(slide, "Product Overview", Inches(0.6), Inches(0.3), Inches(8),
             font_size=Pt(30), bold=True, color=WHITE)
    add_text(slide, "The operator portal — your command center for every customer interaction",
             Inches(0.6), Inches(0.82), Inches(11), Inches(0.4),
             font_size=Pt(14), color=D300, italic=True)
    # Left column — feature list
    features = [
        ("📊", "Dashboard", "Live metrics on messages, auto-reply rates, escalations, and response times"),
        ("💬", "Messages", "Unified message history across WhatsApp and Email with full thread view"),
        ("🚨", "Escalations", "Intelligent queue with priority ranking and assignment tracking"),
        ("🤖", "Auto-Replies", "Log of AI-generated replies with confidence scores and feedback"),
        ("👥", "Backup Contacts", "Escalation contact hierarchy and fallback routing rules"),
        ("📄", "Document Hub", "Knowledge base linked to AI RAG for accurate, context-aware replies"),
        ("🔑", "Key Phrase Maps", "Custom trigger phrases mapped to specific AI reply strategies"),
        ("📱", "Channels", "WhatsApp and Email configuration, webhook status, and connection health"),
        ("💳", "Subscription", "Plan management, usage tracking, and billing history"),
    ]
    fx = Inches(0.5)
    fy = Inches(1.55)
    fw = Inches(4.8)
    for icon, title, desc in features:
        add_rect(slide, fx, fy, Inches(0.45), Inches(0.45), fill_rgb=D100)
        add_text(slide, icon, fx, fy, Inches(0.45), Inches(0.45),
                 font_size=Pt(14), align=PP_ALIGN.CENTER)
        add_text(slide, title, fx + Inches(0.55), fy,
                 Inches(1.2), Inches(0.3), font_size=Pt(12), bold=True, color=D800)
        add_text(slide, desc, fx + Inches(0.55), fy + Inches(0.28),
                 fw - Inches(0.6), Inches(0.3), font_size=Pt(10), color=SLATE, wrap=True)
        add_rect(slide, fx, fy + Inches(0.58), fw, Inches(0.01), fill_rgb=D200)
        fy += Inches(0.62)
    # Right column — mock dashboard preview
    add_rect(slide, Inches(5.8), Inches(1.55), Inches(7), Inches(5.6),
             fill_rgb=WHITE, line_rgb=D300, line_width=Pt(1.5))
    add_rect(slide, Inches(5.8), Inches(1.55), Inches(7), Inches(0.55), fill_rgb=D800)
    add_text(slide, "  Sequor  Dashboard", Inches(5.8), Inches(1.6),
             Inches(7), Inches(0.45), font_size=Pt(13), bold=True, color=WHITE)
    # Stat cards mock
    stats = [("Messages", "1,284", "+12%", D500), ("Auto-Replied", "1,027", "80%", GREEN),
             ("Escalated", "42", "-3%", AMBER), ("Avg Time", "3.2hr", "4hr SLA", D600)]
    sx = Inches(6.0)
    sy = Inches(2.3)
    for label, val, change, col in stats:
        add_rect(slide, sx, sy, Inches(1.5), Inches(0.85), fill_rgb=D50, line_rgb=D200, line_width=Pt(0.75))
        add_rect(slide, sx, sy, Inches(1.5), Inches(0.04), fill_rgb=col)
        add_text(slide, label, sx + Inches(0.1), sy + Inches(0.1),
                 Inches(1.3), Inches(0.3), font_size=Pt(9), color=SLATE)
        add_text(slide, val, sx + Inches(0.1), sy + Inches(0.35),
                 Inches(1.0), Inches(0.4), font_size=Pt(18), bold=True, color=D800)
        sx += Inches(1.58)
    # Bar chart mock
    sy2 = Inches(3.3)
    add_rect(slide, Inches(6.0), sy2, Inches(3.5), Inches(1.5), fill_rgb=D50, line_rgb=D200, line_width=Pt(0.75))
    add_text(slide, "Messages This Week", Inches(6.1), sy2 + Inches(0.05),
             Inches(2), Inches(0.3), font_size=Pt(9), bold=True, color=D700)
    bars = [(0.4, D400), (0.65, D400), (0.5, D400), (0.8, D500), (0.6, D400), (0.9, D500), (0.7, D400)]
    days = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
    bx = Inches(6.2)
    bw = Inches(0.35)
    for (frac, col), day in zip(bars, days):
        bh = frac * Inches(0.9)
        add_rect(slide, bx, sy2 + Inches(1.0) - bh, bw, bh, fill_rgb=col)
        add_text(slide, day, bx - Inches(0.05), sy2 + Inches(1.02),
                 bw + Inches(0.1), Inches(0.2), font_size=Pt(7), color=SLATE, align=PP_ALIGN.CENTER)
        bx += bw + Inches(0.07)
    # Recent messages mock
    add_rect(slide, Inches(9.8), sy2, Inches(2.9), Inches(1.5), fill_rgb=D50, line_rgb=D200, line_width=Pt(0.75))
    add_text(slide, "Recent Messages", Inches(9.9), sy2 + Inches(0.05),
             Inches(2), Inches(0.3), font_size=Pt(9), bold=True, color=D700)
    msgs = [("📧 alice@co.com", "Order #12345 inquiry", "Auto", GREEN),
            ("💬 +1 555 0000", "Refund request", "Escalated", AMBER),
            ("📧 bob@store.io", "Shipping status", "Auto", GREEN)]
    my = sy2 + Inches(0.35)
    for ch, txt, status, col in msgs:
        add_text(slide, ch, Inches(9.9), my, Inches(2.7), Inches(0.25),
                 font_size=Pt(8.5), color=D700)
        add_text(slide, txt, Inches(9.9), my + Inches(0.2), Inches(2.0), Inches(0.22),
                 font_size=Pt(8), color=SLATE)
        add_rect(slide, Inches(11.85), my + Inches(0.22), Inches(0.65), Inches(0.2),
                 fill_rgb=col)
        add_text(slide, status, Inches(11.85), my + Inches(0.22),
                 Inches(0.65), Inches(0.2), font_size=Pt(7.5), bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        my += Inches(0.38)
    # Sidebar mock
    add_rect(slide, Inches(5.8), Inches(4.95), Inches(1.5), Inches(2.2),
             fill_rgb=WHITE, line_rgb=D200, line_width=Pt(0.75))
    nav_items = ["Dashboard", "Messages", "Escalations", "Auto-Replies",
                 "Backup Contacts", "Document Hub"]
    ny = Inches(5.1)
    for item in nav_items:
        add_text(slide, item, Inches(5.95), ny, Inches(1.2), Inches(0.3),
                 font_size=Pt(9), color=D700)
        ny += Inches(0.33)


# ── SLIDE 5: AI Pipeline ───────────────────────────────────────────────────────
def slide_ai():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=D50)
    add_rect(slide, 0, 0, W, Inches(1.3), fill_rgb=D800)
    add_text(slide, "How It Works: AI Pipeline", Inches(0.6), Inches(0.3), Inches(9),
             font_size=Pt(30), bold=True, color=WHITE)
    add_text(slide, "From inbound message to intelligent reply — in seconds",
             Inches(0.6), Inches(0.82), Inches(10), Inches(0.4),
             font_size=Pt(14), color=D300, italic=True)
    steps = [
        (D500, "1", "Inbound",
         "Message arrives via\nWhatsApp or Email.\nWebhook captures\nmetadata + body."),
        (D400, "2", "Classify",
         "MessageClassifier\nidentifies intent:\nrefund, order status,\ncomplaint, or other."),
        (D300, "3", "Retrieve",
         "RAG pipeline searches\nthe Document Hub for\nrelevant knowledge\nbase entries."),
        (D400, "4", "Decide",
         "AutoReplyService\ndetermines: auto-reply\nor escalate to\nhuman operator."),
        (D500, "5", "Respond",
         "AI-generated reply sent\nvia SendGrid or\nWhatsApp API with\naudit trail logged."),
    ]
    sw = Inches(2.2)
    sx = Inches(0.5)
    sy = Inches(2.0)
    arrow_x = sx + sw + Inches(0.05)
    for i, (col, num, title, desc) in enumerate(steps):
        # Box
        add_rect(slide, sx, sy, sw, Inches(3.2), fill_rgb=col)
        # Number circle
        add_rect(slide, sx + Inches(0.85), sy + Inches(0.15),
                 Inches(0.5), Inches(0.5), fill_rgb=WHITE)
        add_text(slide, num, sx + Inches(0.85), sy + Inches(0.15),
                 Inches(0.5), Inches(0.5), font_size=Pt(16), bold=True,
                 color=col, align=PP_ALIGN.CENTER)
        # Title
        add_text(slide, title, sx, sy + Inches(0.75), sw, Inches(0.45),
                 font_size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Description
        add_text(slide, desc, sx + Inches(0.15), sy + Inches(1.3),
                 sw - Inches(0.3), Inches(1.8), font_size=Pt(12),
                 color=WHITE, align=PP_ALIGN.CENTER, wrap=True)
        # Arrow (except last)
        if i < len(steps) - 1:
            add_text(slide, "→", arrow_x, sy + Inches(1.3), Inches(0.4), Inches(0.5),
                     font_size=Pt(28), bold=True, color=D700, align=PP_ALIGN.CENTER)
        sx += sw + Inches(0.5)
    # Learning loop annotation
    add_rect(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2),
             fill_rgb=WHITE, line_rgb=D200, line_width=Pt(1))
    add_rect(slide, Inches(0.5), Inches(5.5), Inches(0.06), Inches(1.2), fill_rgb=GREEN)
    add_text(slide, "🔄  Continuous Learning Loop",
             Inches(0.8), Inches(5.6), Inches(4), Inches(0.4),
             font_size=Pt(13), bold=True, color=D800)
    add_text(slide, "After every escalation resolution, the LearningLoop records operator feedback — improving classification accuracy and reply quality over time. Each model update compounds across all future messages.",
             Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.6),
             font_size=Pt(11.5), color=SLATE, wrap=True)


# ── SLIDE 6: WhatsApp + Email ─────────────────────────────────────────────────
def slide_channels():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=D50)
    add_rect(slide, 0, 0, W, Inches(1.3), fill_rgb=D800)
    add_text(slide, "Multi-Channel: WhatsApp + Email", Inches(0.6), Inches(0.3), Inches(9),
             font_size=Pt(30), bold=True, color=WHITE)
    add_text(slide, "Both channels share the same AI pipeline, knowledge base, and operator portal",
             Inches(0.6), Inches(0.82), Inches(11), Inches(0.4),
             font_size=Pt(14), color=D300, italic=True)
    # WhatsApp card
    add_rect(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(4.7),
             fill_rgb=WHITE, line_rgb=D200, line_width=Pt(1))
    add_rect(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.65), fill_rgb=D600)
    add_text(slide, "💬  WhatsApp", Inches(0.7), Inches(1.72), Inches(5), Inches(0.45),
             font_size=Pt(18), bold=True, color=WHITE)
    whatsapp_features = [
        "Meta Cloud API integration — no legacy SDKs",
        "Webhook verification with HMAC signatures",
        "Template message support for outbound",
        "Contact profiling and phone number validation",
        "Media support: images, documents, audio",
        "Delivery and read receipt tracking",
        "Opt-out / unsubscribe handling",
    ]
    wy = Inches(2.45)
    for feat in whatsapp_features:
        add_text(slide, f"✓  {feat}", Inches(0.8), wy, Inches(5.3), Inches(0.38),
                 font_size=Pt(12), color=D700)
        wy += Inches(0.48)
    # Email card
    add_rect(slide, Inches(6.9), Inches(1.6), Inches(5.9), Inches(4.7),
             fill_rgb=WHITE, line_rgb=D200, line_width=Pt(1))
    add_rect(slide, Inches(6.9), Inches(1.6), Inches(5.9), Inches(0.65), fill_rgb=D500)
    add_text(slide, "📧  Email", Inches(7.1), Inches(1.72), Inches(5.5), Inches(0.45),
             font_size=Pt(18), bold=True, color=WHITE)
    email_features = [
        "SendGrid Inbound Parse webhook — full email receipt",
        "Attachment ingestion and Document Hub filing",
        "HTML and plain-text body parsing",
        "Reply-to threading (in_reply_to header)",
        "SPF/DKIM/DMARC aware routing",
        "Auto-reply with quoted original",
        "PDPA: PII minimization in email storage",
    ]
    ey = Inches(2.45)
    for feat in email_features:
        add_text(slide, f"✓  {feat}", Inches(7.2), ey, Inches(5.4), Inches(0.38),
                 font_size=Pt(12), color=D700)
        ey += Inches(0.48)
    # Bottom note
    add_text(slide, "Both channels use the same AI pipeline — classification → RAG → decision → reply",
             Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.4),
             font_size=Pt(12), bold=True, color=D600, align=PP_ALIGN.CENTER)


# ── SLIDE 7: Escalations ────────────────────────────────────────────────────────
def slide_escalations():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=D50)
    add_rect(slide, 0, 0, W, Inches(1.3), fill_rgb=D800)
    add_text(slide, "Smart Escalations", Inches(0.6), Inches(0.3), Inches(8),
             font_size=Pt(30), bold=True, color=WHITE)
    add_text(slide, "The right message, to the right operator, with full context — instantly",
             Inches(0.6), Inches(0.82), Inches(11), Inches(0.4),
             font_size=Pt(14), color=D300, italic=True)
    # Left: trigger types
    add_text(slide, "Escalation Triggers", Inches(0.5), Inches(1.55), Inches(5.5),
             font_size=Pt(16), bold=True, color=D800)
    triggers = [
        (RED,   "Cancellation / Refund Threats", "Message contains high-friction intent signals"),
        (AMBER, "VIP / Priority Contacts",        "Contact flagged as high-value or flagged for escalation"),
        (AMBER, "SLA Breach Risk",                "Message unresolved beyond configured SLA window"),
        (RED,   "Sentiment: Frustration / Anger",  "Classifier confidence > 0.8 for negative sentiment"),
        (RED,   "Repeat Queries (3x+)",            "Same topic sent 3+ times without resolution"),
        (D500,  "Manual Flag",                      "Operator manually escalates from the portal"),
    ]
    ty = Inches(2.05)
    for col, title, desc in triggers:
        add_rect(slide, Inches(0.5), ty, Inches(0.18), Inches(0.45), fill_rgb=col)
        add_text(slide, title, Inches(0.8), ty, Inches(2.5), Inches(0.3),
                 font_size=Pt(12), bold=True, color=D800)
        add_text(slide, desc, Inches(3.3), ty, Inches(3.0), Inches(0.35),
                 font_size=Pt(11), color=SLATE, italic=True)
        ty += Inches(0.6)
    # Right: escalation flow
    add_rect(slide, Inches(7.0), Inches(1.55), Inches(5.8), Inches(5.1),
             fill_rgb=WHITE, line_rgb=D200, line_width=Pt(1))
    add_rect(slide, Inches(7.0), Inches(1.55), Inches(5.8), Inches(0.55), fill_rgb=D700)
    add_text(slide, "  Escalation Resolution Flow", Inches(7.0), Inches(1.62),
             Inches(5.8), Inches(0.4), font_size=Pt(13), bold=True, color=WHITE)
    flow_steps = [
        ("🚨", "Alert Created", "Escalation record created with full message context, contact info, and SLA timer."),
        ("👤", "Assigned", "Auto-assigned to available operator based on team routing rules and contact tier."),
        ("💬", "Context Shared", "Portal shows full message thread + relevant document excerpts from RAG."),
        ("✅", "Resolved", "Operator resolves with outcome logged. LearningLoop records resolution for training."),
        ("📊", "SLA Closed", "Ticket closed. Resolution time, outcome, and satisfaction score stored in audit trail."),
    ]
    fy = Inches(2.3)
    for icon, title, desc in flow_steps:
        add_rect(slide, Inches(7.2), fy, Inches(0.5), Inches(0.5), fill_rgb=D100)
        add_text(slide, icon, Inches(7.2), fy, Inches(0.5), Inches(0.5),
                 font_size=Pt(14), align=PP_ALIGN.CENTER)
        add_text(slide, title, Inches(7.85), fy + Inches(0.02), Inches(2.5), Inches(0.3),
                 font_size=Pt(12), bold=True, color=D800)
        add_text(slide, desc, Inches(7.85), fy + Inches(0.3), Inches(4.6), Inches(0.38),
                 font_size=Pt(10.5), color=SLATE, wrap=True)
        if icon != "📊":
            add_rect(slide, Inches(7.42), fy + Inches(0.55), Inches(0.05), Inches(0.25),
                     fill_rgb=D300)
        fy += Inches(0.82)


# ── SLIDE 8: PDPA Compliance ───────────────────────────────────────────────────
def slide_compliance():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=D50)
    add_rect(slide, 0, 0, W, Inches(1.3), fill_rgb=D800)
    add_text(slide, "PDPA Compliance by Design", Inches(0.6), Inches(0.3), Inches(9),
             font_size=Pt(30), bold=True, color=WHITE)
    add_text(slide, "Every layer of Sequor is built with Singapore data protection requirements in mind",
             Inches(0.6), Inches(0.82), Inches(11), Inches(0.4),
             font_size=Pt(14), color=D300, italic=True)
    pillars = [
        (D500, "🔐", "Data Minimization",
         "Sequor stores only what's operationally necessary. Email bodies and WhatsApp content are processed in-memory and retained only for active threads. Historical archives are pseudonymized."),
        (D600, "🗑️", "Right to Erasure",
         "Full deletion workflow: when a data subject requests erasure, Sequor purges their PII from all active records and audit logs. Erasure confirmation issued automatically."),
        (D700, "📋", "Audit Trail",
         "Every message, classification decision, escalation, and operator action is logged immutably. Logs include tenant_id, timestamp, actor, and outcome — PDPA Article 20 compliant."),
        (D500, "🔎", "Access Control",
         "Role-based access: operators see only their assigned contacts. Tenant isolation enforced at the database query layer. No cross-tenant data leakage possible."),
        (D600, "📝", "Consent Management",
         "Inbound contacts are flagged for consent status. Audit log records when consent was given, for which channel, and under which privacy notice version."),
        (D700, "🏷️", "Data Classification",
         "PII fields (email addresses, phone numbers, names) are classified and labeled in the audit trail. RAG retrieval respects classification — no PII in AI training data."),
    ]
    pw = Inches(3.9)
    px = Inches(0.5)
    py = Inches(1.6)
    for i, (col, icon, title, desc) in enumerate(pillars):
        row = i // 3
        col_idx = i % 3
        x = px + col_idx * (pw + Inches(0.22))
        y = py + row * Inches(2.5)
        add_rect(slide, x, y, pw, Inches(2.2), fill_rgb=WHITE, line_rgb=D200, line_width=Pt(1))
        add_rect(slide, x, y, pw, Inches(0.06), fill_rgb=col)
        add_text(slide, icon + "  " + title, x + Inches(0.2), y + Inches(0.18),
                 pw - Inches(0.4), Inches(0.45), font_size=Pt(14), bold=True, color=D800)
        add_text(slide, desc, x + Inches(0.2), y + Inches(0.68),
                 pw - Inches(0.4), Inches(1.4), font_size=Pt(11), color=SLATE, wrap=True)


# ── SLIDE 9: Pricing ───────────────────────────────────────────────────────────
def slide_pricing():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=D50)
    add_rect(slide, 0, 0, W, Inches(1.3), fill_rgb=D800)
    add_text(slide, "Pricing", Inches(0.6), Inches(0.3), Inches(6),
             font_size=Pt(30), bold=True, color=WHITE)
    add_text(slide, "Simple, scalable plans — free forever for development",
             Inches(0.6), Inches(0.82), Inches(10), Inches(0.4),
             font_size=Pt(14), color=D300, italic=True)
    plans = [
        ("Free", "$0", "/mo", [
            "50 messages/mo",
            "WhatsApp + Email",
            "Auto-reply AI",
            "1 operator",
            "7-day message history",
            "Community support",
        ], D200, D700, False),
        ("Solo", "$15", "mo", [
            "200 messages/mo",
            "WhatsApp + Email",
            "Auto-reply AI",
            "RAG Document Hub",
            "30-day message history",
            "1 operator",
            "Email support",
        ], D200, D700, False),
        ("Starter", "$35", "mo", [
            "Unlimited messages*",
            "WhatsApp + Email",
            "Auto-reply AI + RAG",
            "Smart Escalations",
            "PDPA audit trail",
            "90-day message history",
            "Up to 3 operators",
            "Priority support",
        ], D500, WHITE, True),
        ("Professional", "$55", "mo", [
            "Everything in Starter",
            "Unlimited operators",
            "Unlimited documents",
            "Priority RAG processing",
            "Advanced analytics",
            "Custom key phrase maps",
            "PDPA compliance report",
            "Dedicated support",
        ], D800, WHITE, False),
    ]
    pw = Inches(2.9)
    px = Inches(0.45)
    py = Inches(1.6)
    for name, price, period, features, col, txt_col, featured in plans:
        add_rect(slide, px, py, pw, Inches(5.2),
                 fill_rgb=col, line_rgb=col if featured else D200,
                 line_width=Pt(2) if featured else Pt(1))
        if featured:
            add_rect(slide, px, py, pw, Inches(0.45), fill_rgb=D400)
            add_text(slide, "★ Most Popular", px, py + Inches(0.07), pw, Inches(0.35),
                     font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_text(slide, name, px, py + Inches(0.55), pw, Inches(0.5),
                     font_size=Pt(20), bold=True, color=txt_col, align=PP_ALIGN.CENTER)
        else:
            add_text(slide, name, px, py + Inches(0.2), pw, Inches(0.5),
                     font_size=Pt(20), bold=True, color=txt_col, align=PP_ALIGN.CENTER)
        add_text(slide, price, px, py + (Inches(0.8) if featured else Inches(0.7)),
                 pw, Inches(0.8), font_size=Pt(38), bold=True,
                 color=txt_col, align=PP_ALIGN.CENTER)
        add_text(slide, period, px, py + (Inches(1.5) if featured else Inches(1.4)),
                 pw, Inches(0.3), font_size=Pt(12), color=txt_col,
                 align=PP_ALIGN.CENTER)
        add_rect(slide, px + Inches(0.3), py + (Inches(1.9) if featured else Inches(1.8)),
                 pw - Inches(0.6), Inches(0.02), fill_rgb=txt_col)
        fy = py + (Inches(2.05) if featured else Inches(1.95))
        for feat in features:
            add_text(slide, f"✓  {feat}", px + Inches(0.25), fy,
                     pw - Inches(0.4), Inches(0.38), font_size=Pt(11),
                     color=txt_col)
            fy += Inches(0.42)
        px += pw + Inches(0.22)
    add_text(slide, "* Message limits enforced at fair-use threshold. Enterprise volume pricing available on request.",
             Inches(0.5), Inches(6.85), Inches(12), Inches(0.3),
             font_size=Pt(10), color=SLATE, italic=True)


# ── SLIDE 10: Live Demo ────────────────────────────────────────────────────────
def slide_demo():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=D800)
    # Big decorative shapes
    add_rect(slide, Inches(8), Inches(-2), Inches(6), Inches(6), fill_rgb=D700)
    add_rect(slide, Inches(9.5), Inches(-1), Inches(5), Inches(5), fill_rgb=D600)
    add_rect(slide, Inches(-2), Inches(4.5), Inches(5), Inches(5), fill_rgb=D700)
    add_rect(slide, Inches(-1), Inches(5), Inches(4), Inches(4), fill_rgb=D600)
    add_text(slide, "Live Demo", Inches(0.8), Inches(1.0), Inches(8),
             font_size=Pt(48), bold=True, color=WHITE)
    add_rect(slide, Inches(0.8), Inches(1.75), Inches(4), Inches(0.06), fill_rgb=D500)
    add_text(slide, "Watch Sequor handle a real inbound message end-to-end",
             Inches(0.8), Inches(1.95), Inches(8), Inches(0.5),
             font_size=Pt(18), color=D300)
    # Demo flow steps
    steps = [
        ("Step 1", "Inbound message received via WhatsApp or Email webhook"),
        ("Step 2", "AI classifies intent and searches the knowledge base (RAG)"),
        ("Step 3", "Auto-reply generated and sent — or escalation triggered"),
        ("Step 4", "Portal dashboard and message history update in real time"),
        ("Step 5", "Escalation queue shows new item with full context"),
    ]
    sx = Inches(0.8)
    sy = Inches(2.8)
    for step, desc in steps:
        add_rect(slide, sx, sy, Inches(0.08), Inches(0.55), fill_rgb=D500)
        add_text(slide, step, sx + Inches(0.2), sy, Inches(1.2), Inches(0.3),
                 font_size=Pt(11), bold=True, color=D400)
        add_text(slide, desc, sx + Inches(0.2), sy + Inches(0.28),
                 Inches(7), Inches(0.4), font_size=Pt(13), color=WHITE)
        sy += Inches(0.72)
    add_text(slide, "Try it now — scan the QR or visit the demo portal live",
             Inches(0.8), Inches(6.5), Inches(8), Inches(0.4),
             font_size=Pt(14), bold=True, color=D300)


# ── SLIDE 11: Roadmap ──────────────────────────────────────────────────────────
def slide_roadmap():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=D50)
    add_rect(slide, 0, 0, W, Inches(1.3), fill_rgb=D800)
    add_text(slide, "Roadmap", Inches(0.6), Inches(0.3), Inches(8),
             font_size=Pt(30), bold=True, color=WHITE)
    add_text(slide, "What's coming next for Sequor",
             Inches(0.6), Inches(0.82), Inches(8), Inches(0.4),
             font_size=Pt(14), color=D300, italic=True)
    quarters = [
        ("Q3 2026", D500, [
            ("SMS Channel", "Twilio integration for SMS inbound/outbound"),
            ("Team Roles", "Admin, Manager, Operator access tiers"),
            ("Analytics Dashboard", "Message volume trends, AI accuracy reports, SLA dashboards"),
        ]),
        ("Q4 2026", D600, [
            ("Multi-language AI", "AI replies in Thai, Bahasa, Vietnamese, Mandarin"),
            ("API Webhooks", "Custom outbound webhook triggers for CRM/ERP integration"),
            ("Mobile App", "iOS/Android operator app for escalation management on the go"),
        ]),
        ("Q1 2027", D700, [
            ("Enterprise SSO", "SAML/OIDC single sign-on for enterprise customers"),
            ("SOC 2 Type II", "Security certification for enterprise procurement"),
            ("Voice Channel", "Inbound voice with transcription and AI summarization"),
        ]),
    ]
    rx = Inches(0.5)
    rw = Inches(3.9)
    ry = Inches(1.6)
    for qtr, col, items in quarters:
        add_rect(slide, rx, ry, rw, Inches(5.3), fill_rgb=WHITE, line_rgb=D200, line_width=Pt(1))
        add_rect(slide, rx, ry, rw, Inches(0.65), fill_rgb=col)
        add_text(slide, qtr, rx, ry + Inches(0.15), rw, Inches(0.45),
                 font_size=Pt(17), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        iy = ry + Inches(0.85)
        for title, desc in items:
            add_rect(slide, rx + Inches(0.2), iy, Inches(0.06), Inches(0.45), fill_rgb=col)
            add_text(slide, title, rx + Inches(0.35), iy,
                     rw - Inches(0.5), Inches(0.32), font_size=Pt(13), bold=True, color=D800)
            add_text(slide, desc, rx + Inches(0.35), iy + Inches(0.33),
                     rw - Inches(0.5), Inches(0.6), font_size=Pt(11), color=SLATE, wrap=True)
            iy += Inches(1.15)
        rx += rw + Inches(0.28)


# ── SLIDE 12: Why Sequor ───────────────────────────────────────────────────────
def slide_why():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=D50)
    add_rect(slide, 0, 0, W, Inches(1.3), fill_rgb=D800)
    add_text(slide, "Why Sequor", Inches(0.6), Inches(0.3), Inches(8),
             font_size=Pt(30), bold=True, color=WHITE)
    add_text(slide, "Built for Southeast Asian SMBs who need enterprise-grade AI without the complexity",
             Inches(0.6), Inches(0.82), Inches(12), Inches(0.4),
             font_size=Pt(14), color=D300, italic=True)
    reasons = [
        (D500, "Built for Southeast Asia",
         "Multi-language AI, WhatsApp as primary channel, PDPA-first compliance — designed for the region's SMB reality, not adapted from Western enterprise software."),
        (D600, "10x Operator Efficiency",
         "80% of messages auto-replied. Operators spend time on the 20% that matters — escalations, VIPs, and complex cases — not repetitive FAQs."),
        (D700, "Zero Infrastructure Headaches",
         "No setup, no servers, no ML pipelines. Connect your WhatsApp Business and SendGrid accounts in 2 minutes. The AI learns from your documents."),
        (D500, "PDPA from Day One",
         "Compliance isn't an afterthought or an add-on module. Audit trails, erasure workflows, and data minimization are built into every layer."),
        (D600, "Transparent Pricing",
         "Predictable per-seat pricing. No per-message AI credits, no hidden ingestion fees, no enterprise-only features locked behind a sales call."),
        (D700, "Fast Time to Value",
         "Customers see their first AI reply within 10 minutes of signing up. Knowledge base built from their existing documents — no training data required."),
    ]
    rw = Inches(3.9)
    rx = Inches(0.5)
    ry = Inches(1.6)
    for i, (col, title, desc) in enumerate(reasons):
        row = i // 3
        col_idx = i % 3
        x = rx + col_idx * (rw + Inches(0.22))
        y = ry + row * Inches(2.5)
        add_rect(slide, x, y, rw, Inches(2.2), fill_rgb=WHITE, line_rgb=D200, line_width=Pt(1))
        add_rect(slide, x, y, rw, Inches(0.06), fill_rgb=col)
        add_text(slide, title, x + Inches(0.2), y + Inches(0.2),
                 rw - Inches(0.4), Inches(0.5), font_size=Pt(14), bold=True, color=D800)
        add_text(slide, desc, x + Inches(0.2), y + Inches(0.75),
                 rw - Inches(0.4), Inches(1.3), font_size=Pt(11.5), color=SLATE, wrap=True)


# ── SLIDE 13: The Ask ─────────────────────────────────────────────────────────
def slide_ask():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=D800)
    add_rect(slide, Inches(8), Inches(-1), Inches(6), Inches(6), fill_rgb=D700)
    add_rect(slide, Inches(9.5), Inches(-0.5), Inches(5), Inches(5), fill_rgb=D600)
    add_rect(slide, Inches(-1.5), Inches(5), Inches(4), Inches(4), fill_rgb=D700)
    add_text(slide, "The Opportunity", Inches(0.8), Inches(0.8), Inches(8),
             font_size=Pt(40), bold=True, color=WHITE)
    add_rect(slide, Inches(0.8), Inches(1.6), Inches(4), Inches(0.05), fill_rgb=D500)
    stats = [
        ("$2.4B", "TAM: SE Asian SMB Customer\nCommunication Software (2026)"),
        ("35%", "Annual growth in WhatsApp\nBusiness API adoption"),
        ("80%", "Average auto-reply rate\nin Sequor pilot accounts"),
        ("< 4 weeks", "Average customer time-to-first\nAI-replied message"),
    ]
    sx = Inches(0.8)
    sy = Inches(2.0)
    for val, desc in stats:
        add_text(slide, val, sx, sy, Inches(2.5), Inches(0.9),
                 font_size=Pt(36), bold=True, color=D500)
        add_text(slide, desc, sx, sy + Inches(0.85), Inches(2.8), Inches(0.8),
                 font_size=Pt(11.5), color=D300, wrap=True)
        sy += Inches(1.85)
    # Investment ask
    add_rect(slide, Inches(7.5), Inches(1.8), Inches(5.3), Inches(4.5),
             fill_rgb=D700)
    add_text(slide, "Investment Summary", Inches(7.7), Inches(2.0), Inches(4.9), Inches(0.5),
             font_size=Pt(18), bold=True, color=WHITE)
    add_rect(slide, Inches(7.7), Inches(2.55), Inches(4.9), Inches(0.03), fill_rgb=D500)
    rows = [
        ("Raise", "Seed Round — S$1.2M"),
        ("Use of Funds", "Product engineering (40%)\nSales & marketing (35%)\nOperations (25%)"),
        ("Key Milestones",
         "20 paying customers\n1,000 monthly messages processed\n3 channel integrations live"),
        ("Target Close", "Q3 2026"),
    ]
    ry2 = Inches(2.75)
    for label, val in rows:
        add_text(slide, label + ":", Inches(7.7), ry2, Inches(1.6), Inches(0.35),
                 font_size=Pt(11), bold=True, color=D400)
        add_text(slide, val, Inches(9.4), ry2, Inches(3.3), Inches(0.6),
                 font_size=Pt(11.5), color=WHITE, wrap=True)
        ry2 += Inches(0.75)
    add_rect(slide, 0, Inches(6.7), W, Inches(0.8), fill_rgb=D700)
    add_text(slide, "Confidential  •  May 2026  •  Sequor Pte. Ltd.",
             Inches(0.8), Inches(6.78), Inches(6), Inches(0.5),
             font_size=Pt(11), color=D400, italic=True)


# ── SLIDE 14: Contact ──────────────────────────────────────────────────────────
def slide_contact():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, H, fill_rgb=D800)
    add_rect(slide, Inches(9), Inches(-1.5), Inches(5), Inches(5), fill_rgb=D700)
    add_rect(slide, Inches(10.5), Inches(-0.5), Inches(4), Inches(4), fill_rgb=D600)
    add_rect(slide, Inches(-1), Inches(5.5), Inches(3), Inches(3), fill_rgb=D700)
    add_rect(slide, Inches(0.8), Inches(1.8), Inches(0.9), Inches(0.9), fill_rgb=D500)
    add_text(slide, "S", Inches(0.8), Inches(1.82), Inches(0.9), Inches(0.9),
             font_size=Pt(40), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Sequor", Inches(1.9), Inches(1.85), Inches(4), Inches(0.8),
             font_size=Pt(40), bold=True, color=WHITE)
    add_text(slide, "AI-Powered Customer Communication Platform",
             Inches(0.8), Inches(2.85), Inches(8), Inches(0.4),
             font_size=Pt(16), color=D300)
    add_rect(slide, Inches(0.8), Inches(3.35), Inches(5), Inches(0.05), fill_rgb=D500)
    contact_items = [
        ("🌐", "Website", "www.sequor.com"),
        ("📧", "Email", "hello@sequor.com"),
        ("📱", "Phone", "+65 6000 0000"),
        ("💬", "WhatsApp", "+65 9000 0000"),
        ("📍", "Address", "1 Raffles Place, #12-00\nSingapore 048616"),
    ]
    cy = Inches(3.65)
    for icon, label, val in contact_items:
        add_text(slide, icon + "  " + label + ":", Inches(0.8), cy,
                 Inches(2.2), Inches(0.35), font_size=Pt(13), bold=True, color=D400)
        add_text(slide, val, Inches(3.0), cy, Inches(4), Inches(0.45),
                 font_size=Pt(13), color=WHITE)
        cy += Inches(0.55)
    # Right side — thank you
    add_text(slide, "Thank you\nfor your time.", Inches(7.5), Inches(2.5),
             Inches(5), Inches(2), font_size=Pt(38), bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, wrap=True)
    add_rect(slide, Inches(7.5), Inches(4.6), Inches(4.5), Inches(0.05), fill_rgb=D500)
    add_text(slide, "We look forward to building the future of\ncustomer communication — together.",
             Inches(7.5), Inches(4.8), Inches(5), Inches(0.9),
             font_size=Pt(13), color=D300, align=PP_ALIGN.CENTER, wrap=True)
    add_rect(slide, 0, Inches(6.7), W, Inches(0.8), fill_rgb=D700)
    add_text(slide, "Confidential  •  May 2026  •  Sequor Pte. Ltd.",
             Inches(0.8), Inches(6.78), Inches(6), Inches(0.5),
             font_size=Pt(11), color=D400, italic=True)


# ── Generate ───────────────────────────────────────────────────────────────────
slide_cover()
slide_problem()
slide_solution()
slide_product()
slide_ai()
slide_channels()
slide_escalations()
slide_compliance()
slide_pricing()
slide_demo()
slide_roadmap()
slide_why()
slide_ask()
slide_contact()

out = "/Users/aliciapang/Documents/GitHub/Sequor/workspaces/_template/Sequor_Investor_Deck_May2026.pptx"
prs.save(out)
print(f"Saved: {out}")
