"""Generate Sequor PE Investor Deck — White + Light Teal, Clean & Minimal."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

FONT = "Calibri"

# ── Palette: White dominant, light teal accents, dark teal for text ────────────
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG      = RGBColor(0xFF, 0xFF, 0xFF)   # pure white slides
BG_S    = RGBColor(0xF4, 0xFB, 0xFD)   # very subtle section tint
D050    = RGBColor(0xEB, 0xF6, 0xFA)   # lightest tint — large areas
D100    = RGBColor(0xD4, 0xEA, 0xF4)   # card borders
D200    = RGBColor(0xAA, 0xD5, 0xE8)   # dividers, subtle bars
D300    = RGBColor(0x7F, 0xBE, 0xD4)   # muted accents
D400    = RGBColor(0x54, 0xA8, 0xC0)   # medium accent
D500    = RGBColor(0x29, 0x91, 0xAA)   # primary accent — teal
D600    = RGBColor(0x1E, 0x75, 0x8F)   # supporting accent
D700    = RGBColor(0x17, 0x5A, 0x73)   # dark supporting
D800    = RGBColor(0x14, 0x47, 0x5D)   # darkest — headers, primary text
BODY    = RGBColor(0x4A, 0x5E, 0x6E)   # body text
MUTED   = RGBColor(0x8A, 0x9E, 0xB0)   # muted / labels
GREEN   = RGBColor(0x1A, 0xBE, 0x5C)
AMBER   = RGBColor(0xE8, 0x8A, 0x00)
RED     = RGBColor(0xE0, 0x3A, 0x3A)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def R(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.line.width = lw
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
    else:
        s.line.fill.background()
    return s


def T(slide, text, l, t, w, h,
      size=14, bold=False, italic=False,
      color=BODY, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    return tb


def TCTR(slide, text, l, t, w, h,
         size=14, bold=False, color=BODY):
    """Centered text, vertically centered in box."""
    from pptx.enum.text import MSO_ANCHOR
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    return tb


def HR(slide, l, t, w, color=D200, thick=Pt(1)):
    R(slide, l, t, w, thick, fill=color)


def bg_white(s):
    R(s, 0, 0, W, H, fill=WHITE)


def header_bar(s, title, bg=D800):
    R(s, 0, 0, W, Inches(1.05), fill=bg)
    T(s, title, Inches(0.6), Inches(0.25), Inches(12.0), Inches(0.62),
      size=26, bold=True, color=WHITE)


def card(s, l, t, w, h, fill=WHITE, border=D100, border_w=Pt(0.75)):
    R(s, l, t, w, h, fill=fill, line=border, lw=border_w)


def footer(s):
    R(s, 0, Inches(7.1), W, Pt(1), fill=D100)
    T(s, "Sequor Pte. Ltd.  |  Confidential  |  May 2026",
      Inches(0.6), Inches(7.18), Inches(12.0), Inches(0.3),
      size=10, color=MUTED)


# ── SLIDE 1: Cover ────────────────────────────────────────────────────────────
def s1():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)

    # Thin teal accent bar at very top
    R(s, 0, 0, W, Inches(0.12), fill=D500)

    # Logo mark — small S in teal circle
    R(s, Inches(0.7), Inches(1.5), Inches(1.0), Inches(1.0), fill=D500)
    TCTR(s, "S", Inches(0.7), Inches(1.5), Inches(1.0), Inches(1.0),
         size=42, bold=True, color=WHITE)

    # Company name
    T(s, "Sequor", Inches(2.0), Inches(1.62), Inches(5.0), Inches(0.7),
      size=40, bold=True, color=D800)
    T(s, "Pte. Ltd.", Inches(2.0), Inches(2.3), Inches(5.0), Inches(0.4),
      size=16, color=MUTED)

    # Horizontal rule
    HR(s, Inches(0.7), Inches(3.0), Inches(5.5), color=D200)

    # Main headline
    T(s, "Investor Presentation",
      Inches(0.7), Inches(3.3), Inches(8.0), Inches(0.8),
      size=30, bold=True, color=D800)

    # Subtitle
    T(s, "Southeast Asian SMBs deserve an AI coverage layer\n"
         "that works while they're busy — not just reminds them.",
      Inches(0.7), Inches(4.25), Inches(8.0), Inches(1.0),
      size=16, color=BODY)

    # Date
    T(s, "May 2026", Inches(0.7), Inches(5.5), Inches(3.0), Inches(0.4),
      size=14, color=MUTED)

    # Right side — large teal circle with tagline
    R(s, Inches(9.5), Inches(1.2), Inches(3.5), Inches(3.5), fill=D050)
    R(s, Inches(9.8), Inches(1.5), Inches(2.9), Inches(2.9), fill=D200)
    TCTR(s, "AI Coverage\nfor Business",
         Inches(9.8), Inches(1.8), Inches(2.9), Inches(2.2),
         size=20, bold=True, color=D700)

    footer(s)


# ── SLIDE 2: Problem ─────────────────────────────────────────────────────────
def s2():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    header_bar(s, "The Problem")
    footer(s)

    # Lead stat
    R(s, Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.9), fill=D050)
    T(s, "Southeast Asian SMB teams lose 12–18 hours/week to coverage gaps — "
         "unanswered messages, manual routing, and follow-up that the AI should handle.",
      Inches(0.85), Inches(1.35), Inches(11.6), Inches(0.7),
      size=14, bold=True, color=D700)

    # Four pain points
    problems = [
        ("No coverage when away",
         "Customers message anytime. But your team works 9–6. There's nothing between 'message received' and 'human available.'"),
        ("Every query hits a human first",
         "Routine questions — pricing, availability, policy — consume your team's most expensive hours."),
        ("Tool silos everywhere",
         "Email in one place, WhatsApp in another, tasks in a third. No unified view of what's resolved."),
        ("AI tools need adoption",
         "Existing AI tools require your team to learn new workflows. Sequor is invisible to employees — they just reply to escalations."),
    ]
    px = Inches(0.6)
    pw = Inches(5.95)
    py = Inches(2.35)
    for title, desc in problems:
        R(s, px, py, Inches(0.08), Inches(0.85), fill=D500)
        T(s, title, px + Inches(0.22), py, pw - Inches(0.3), Inches(0.38),
          size=13, bold=True, color=D800)
        T(s, desc, px + Inches(0.22), py + Inches(0.4), pw - Inches(0.3), Inches(0.52),
          size=11, color=BODY, wrap=True)
        py += Inches(1.08)
        if py > Inches(5.8):
            py = Inches(2.35)
            px += pw + Inches(0.28)


# ── SLIDE 3: Solution ─────────────────────────────────────────────────────────
def s3():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    header_bar(s, "Our Solution")
    footer(s)

    T(s, "Sequor is the AI coverage layer for customer communication.",
      Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.6),
      size=20, bold=True, color=D800)
    T(s, "An always-on layer between your team and your customers — handles routine queries, escalates what it can't, learns from every resolution.",
      Inches(0.6), Inches(1.9), Inches(12.1), Inches(0.5),
      size=14, color=BODY)

    HR(s, Inches(0.6), Inches(2.52), Inches(12.1), color=D200)

    # 4 pillars
    pillars = [
        ("Reads Everything",   "WhatsApp + Email in one unified thread. No switching tabs."),
        ("Resolves Routine",   "RAG-powered answers from your documents. No human for routine queries."),
        ("Escalates Smart",    "What it can't answer, routed with full context. PII redacted."),
        ("Learns & Improves",  "Every human correction feeds back. Gets smarter with every interaction."),
    ]
    px = Inches(0.6)
    pw = Inches(2.9)
    py = Inches(2.7)
    for title, desc in pillars:
        R(s, px, py, pw, Inches(0.08), fill=D500)
        T(s, title, px, py + Inches(0.2), pw, Inches(0.42),
          size=14, bold=True, color=D700)
        T(s, desc, px, py + Inches(0.68), pw, Inches(0.85),
          size=12, color=BODY, wrap=True)
        px += pw + Inches(0.25)

    # Bottom tagline
    R(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.85), fill=D050)
    T(s, "Sequor becomes your team's institutional memory — covering when you're away, "
         "remembering what was said, routing what needs humans.",
      Inches(0.8), Inches(6.02), Inches(11.7), Inches(0.6),
      size=13, italic=True, color=D700, align=PP_ALIGN.CENTER)


# ── SLIDE 4: Product ─────────────────────────────────────────────────────────
def s4():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    header_bar(s, "Product: The Portal")
    footer(s)

    # Left — what you see
    T(s, "One view. Everything covered.",
      Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.5),
      size=17, bold=True, color=D800)

    features = [
        "Unified WhatsApp + Email inbox",
        "AI draft responses — one-tap send or edit",
        "Escalation queue with full context",
        "PDPA-compliant audit trail",
        "SLA timers + contact history",
    ]
    fy = Inches(1.95)
    for f in features:
        R(s, Inches(0.6), fy + Inches(0.08), Inches(0.06), Inches(0.38), fill=D500)
        T(s, f, Inches(0.82), fy, Inches(5.4), Inches(0.45),
          size=12, color=BODY)
        fy += Inches(0.55)

    # Right — portal mockup
    card(s, Inches(6.7), Inches(1.25), Inches(6.0), Inches(5.6), border=D100)
    R(s, Inches(6.7), Inches(1.25), Inches(6.0), Inches(0.5), fill=D700)
    T(s, "Portal Preview", Inches(6.9), Inches(1.34), Inches(5.6), Inches(0.35),
      size=12, bold=True, color=WHITE)

    # Channel tabs
    R(s, Inches(6.9), Inches(1.88), Inches(2.2), Inches(0.32), fill=D500)
    T(s, "WhatsApp", Inches(6.9), Inches(1.88), Inches(2.2), Inches(0.32),
      size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(s, "Email  All Contacts", Inches(9.2), Inches(1.88), Inches(3.3), Inches(0.32),
      size=10, color=MUTED)

    # Inbox section
    R(s, Inches(6.9), Inches(2.28), Inches(5.6), Inches(0.65), fill=BG_S)
    T(s, "Inbox — 3 escalations, 12 auto-replied", Inches(7.05), Inches(2.34),
      Inches(5.3), Inches(0.3), size=11, bold=True, color=D700)

    # Message rows
    msgs = [
        ("Customer", "Hi, is the workshop on for Saturday?"),
        ("→ AI Draft", "Hi! Yes, the workshop is confirmed..."),
        ("Customer", "Can I get a refund if I cancel?"),
        ("→ Escalated", "VIP + Refund — routed to operator"),
    ]
    ry = Inches(3.0)
    for sender, msg in msgs:
        sender_col = D700 if sender.startswith("→") else MUTED
        T(s, sender, Inches(6.9), ry, Inches(1.3), Inches(0.28),
          size=9, bold=True, color=sender_col)
        T(s, msg, Inches(8.2), ry, Inches(4.5), Inches(0.28),
          size=10, color=BODY)
        ry += Inches(0.35)

    # Audit bar
    R(s, Inches(6.9), Inches(4.55), Inches(5.6), Inches(0.42), fill=D050)
    T(s, "PDPA Audit  |  12 actions today  |  0 PII leaks  |  All logs immutable",
      Inches(7.05), Inches(4.6), Inches(5.3), Inches(0.32),
      size=10, color=MUTED)

    # Doc hub
    R(s, Inches(6.9), Inches(5.05), Inches(5.6), Inches(1.6), fill=BG_S)
    T(s, "Document Hub", Inches(7.05), Inches(5.12), Inches(5.3), Inches(0.3),
      size=11, bold=True, color=D700)
    T(s, "47 docs indexed  |  Learning Loop: 234 corrections  |  Avg response: 4 min",
      Inches(7.05), Inches(5.45), Inches(5.3), Inches(0.28),
      size=10, color=BODY)
    T(s, "Last AI answer confidence: 91%",
      Inches(7.05), Inches(5.78), Inches(5.3), Inches(0.28),
      size=10, color=GREEN)


# ── SLIDE 5: AI Pipeline ─────────────────────────────────────────────────────
def s5():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    header_bar(s, "AI Pipeline")
    footer(s)

    T(s, "How a customer message becomes a resolution — in minutes.",
      Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.42),
      size=14, color=BODY)

    steps = [
        ("1", "Inbound",         "WhatsApp / Email normalized into a single thread."),
        ("2", "Classify",        "Intent labeled: Routine / Escalation / VIP. Confidence scored."),
        ("3", "Retrieve",        "RAG queries your Document Hub. Top chunks retrieved with citation."),
        ("4", "Draft",           "LLM generates response from chunks. Marked [DRAFT]."),
        ("5", "Send or Hold",    "Routine + High Confidence → auto-send. Anything else → queue."),
        ("6", "Learn",            "Operator edits / approves → feedback logged to Learning Loop."),
    ]
    sw = Inches(2.0)
    sx = Inches(0.5)
    sy = Inches(1.8)
    for num, title, desc in steps:
        # Number bubble
        R(s, sx, sy, Inches(0.5), Inches(0.5), fill=D500)
        TCTR(s, num, sx, sy, Inches(0.5), Inches(0.5),
             size=14, bold=True, color=WHITE)
        # Connector line
        if num != "6":
            R(s, sx + Inches(0.22), sy + Inches(0.52), Inches(0.06), Inches(0.45), fill=D200)
        # Text
        T(s, title, sx + Inches(0.65), sy, Inches(1.3), Inches(0.42),
          size=12, bold=True, color=D700)
        T(s, desc, sx + Inches(0.65), sy + Inches(0.42), Inches(1.3), Inches(0.7),
          size=10, color=BODY, wrap=True)
        sx += sw + Inches(0.18)
        if sx > Inches(10.5):
            sx = Inches(0.5)
            sy += Inches(1.35)


# ── SLIDE 6: Channels ─────────────────────────────────────────────────────────
def s6():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    header_bar(s, "Channels")
    footer(s)

    channels = [
        ("WhatsApp",    D500,
         ["Meta WhatsApp Business API",
          "Company-owned number (not personal)",
          "Template messages for OOO windows",
          "Session management + SLA tracking"]),
        ("Email",        D600,
         ["SendGrid Inbound Parse webhook",
          "Forward or BCC address to portal",
          "Full thread context preserved",
          "PDPA-compliant subject + body"]),
        ("API / Web",   D700,
         ["REST API for custom integrations",
          "Webhook outbound to CRMs",
          "Zapier / Make.com connectors",
          "Enterprise SSO + SCIM"]),
    ]
    cx = Inches(0.6)
    cw = Inches(3.95)
    cy = Inches(1.35)
    ch = Inches(5.4)
    for name, col, bullets in channels:
        card(s, cx, cy, cw, ch, border=D100)
        # Color bar at top of card
        R(s, cx, cy, cw, Inches(0.55), fill=col)
        TCTR(s, name, cx, cy + Inches(0.08), cw, Inches(0.45),
             size=16, bold=True, color=WHITE)
        by = cy + Inches(0.72)
        for b in bullets:
            R(s, cx + Inches(0.2), by + Inches(0.1), Inches(0.06), Inches(0.3), fill=col)
            T(s, b, cx + Inches(0.38), by, cw - Inches(0.56), Inches(0.45),
              size=11, color=BODY, wrap=True)
            by += Inches(0.58)
        cx += cw + Inches(0.28)


# ── SLIDE 7: Competitor Analysis ──────────────────────────────────────────────
def s7():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    header_bar(s, "Competitive Landscape")
    footer(s)

    # Lead line
    T(s, "Sequor is not a better version of existing tools. It's a new category: "
         "the company's communication coverage layer.",
      Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.5),
      size=14, bold=True, color=D700)

    HR(s, Inches(0.6), Inches(1.82), Inches(12.1), color=D100)

    # Competitor vs Sequor comparison
    comparison = [
        ("OpenClaw\nPersonal WhatsApp",
         "Personal WhatsApp automation.\nNo company accounts, no email,\nno routing, no compliance.",
         D500, D700),
        ("Superhuman / SaneBox\nAI Email",
         "Reply faster when present.\nNo coverage when absent.\nNo WhatsApp or routing.",
         D600, D700),
        ("Zendesk / Intercom\nEnterprise Platforms",
         "Manual routing, no autonomous\nresolution. Days to configure.\nNot built for SE Asia SMBs.",
         D600, D700),
        ("Tidio / Freshdesk\nShared Inbox",
         "Basic AI assist.\nNo RAG, no learning loop,\nno PDPA compliance.",
         D700, D700),
    ]
    cw = Inches(2.85)
    cx = Inches(0.6)
    cy = Inches(2.0)
    for name, desc, col, txt in comparison:
        card(s, cx, cy, cw, Inches(2.2), border=D100)
        R(s, cx, cy, cw, Inches(0.08), fill=col)
        T(s, name, cx + Inches(0.15), cy + Inches(0.18), cw - Inches(0.3), Inches(0.72),
          size=11, bold=True, color=col, wrap=True)
        T(s, desc, cx + Inches(0.15), cy + Inches(0.98), cw - Inches(0.3), Inches(1.0),
          size=10, color=BODY, wrap=True)
        cx += cw + Inches(0.22)

    # Sequor highlight bar
    cy = Inches(4.35)
    R(s, Inches(0.6), cy, Inches(12.1), Inches(2.5), fill=D050, line=D500, lw=Pt(1.5))
    R(s, Inches(0.6), cy, Inches(0.12), Inches(2.5), fill=D500)

    T(s, "Sequor's Structural Differentiation",
      Inches(0.9), cy + Inches(0.15), Inches(11.6), Inches(0.45),
      size=15, bold=True, color=D800)

    diffs = [
        ("Company-owned WhatsApp accounts",   "Enterprise-grade, auditable, not tied to personal phone numbers"),
        ("Email-first unified interface",     "WhatsApp + email in one thread — no tool switching"),
        ("Learning loop that compounds",      "Every correction improves future responses — competitors can't shortcut this"),
        ("Smart escalations with PII redaction", "Context delivered, PII protected, SLA tracked end-to-end"),
        ("PDPA compliance built in",          "Audit trail, erasure, consent — not add-ons you pay extra for"),
    ]
    dy = cy + Inches(0.65)
    for label, detail in diffs:
        T(s, f"{label}:", Inches(0.9), dy, Inches(2.8), Inches(0.35),
          size=11, bold=True, color=D700)
        T(s, detail, Inches(3.75), dy, Inches(8.7), Inches(0.35),
          size=11, color=BODY)
        dy += Inches(0.4)


# ── SLIDE 8: Escalations ─────────────────────────────────────────────────────
def s8():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    header_bar(s, "Smart Escalations")
    footer(s)

    # Left — triggers
    T(s, "What triggers an escalation",
      Inches(0.6), Inches(1.3), Inches(5.6), Inches(0.42),
      size=14, bold=True, color=D800)

    triggers = [
        (RED,   "Cancellation / Refund Threats",  "High-friction intent. Confidence > 0.8 → immediate route."),
        (AMBER, "VIP / Priority Contacts",          "Tiered routing — top contacts get priority queue."),
        (AMBER, "SLA Breach Risk",                 "Auto-escalate past your configured SLA window."),
        (RED,   "Negative Sentiment",              "Frustration or anger detected → route regardless of topic."),
        (RED,   "Repeat Queries 3x+",            "Same unanswered topic → escalate to human."),
        (D500,  "Manual Flag",                    "Any thread, one click. Operator flags what the AI missed."),
    ]
    ty = Inches(1.88)
    for col, title, desc in triggers:
        R(s, Inches(0.6), ty + Inches(0.1), Inches(0.08), Inches(0.38), fill=col)
        T(s, title, Inches(0.82), ty, Inches(2.6), Inches(0.38),
          size=12, bold=True, color=D800)
        T(s, desc, Inches(0.82), ty + Inches(0.4), Inches(5.3), Inches(0.38),
          size=11, color=BODY, wrap=True)
        ty += Inches(0.8)

    # Right — flow
    card(s, Inches(6.7), Inches(1.25), Inches(6.0), Inches(5.6), border=D100)
    R(s, Inches(6.7), Inches(1.25), Inches(6.0), Inches(0.55), fill=D700)
    T(s, "Resolution Flow", Inches(6.9), Inches(1.35), Inches(5.6), Inches(0.38),
      size=14, bold=True, color=WHITE)

    flow = [
        ("1", "Alert Created",       "Full context captured. SLA countdown started."),
        ("2", "Assigned",            "Auto-assigned by routing rules + contact tier."),
        ("3", "Context Delivered",   "Thread + RAG excerpts shown in portal. PII redacted."),
        ("4", "Operator Resolves",   "Outcome logged. AI learns from the resolution."),
        ("5", "SLA Closed",          "Resolution time stored. Audit trail immutable."),
    ]
    fy = Inches(2.0)
    for num, title, desc in flow:
        R(s, Inches(6.9), fy, Inches(0.48), Inches(0.48), fill=D500)
        TCTR(s, num, Inches(6.9), fy, Inches(0.48), Inches(0.48),
             size=13, bold=True, color=WHITE)
        T(s, title, Inches(7.55), fy + Inches(0.02), Inches(2.8), Inches(0.32),
          size=12, bold=True, color=D700)
        T(s, desc, Inches(7.55), fy + Inches(0.36), Inches(4.9), Inches(0.38),
          size=11, color=BODY)
        if num != "5":
            R(s, Inches(7.1), fy + Inches(0.55), Inches(0.05), Inches(0.2), fill=D200)
        fy += Inches(0.88)


# ── SLIDE 9: PDPA ─────────────────────────────────────────────────────────────
def s9():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    header_bar(s, "PDPA Compliance by Design")
    footer(s)

    pillars = [
        ("Data Minimization",    "Only operationally necessary data stored. Historical archives pseudonymized."),
        ("Right to Erasure",    "Full deletion workflow. PII purged from records and audit logs."),
        ("Audit Trail",         "Every action logged: tenant, timestamp, actor, outcome. Immutable."),
        ("Tenant Isolation",    "Role-based access. Operators see only their contacts. No leaks."),
        ("Consent Management",   "Contact consent flagged by channel and notice version."),
        ("PII Classification",  "Email, phone, name classified. RAG respects boundaries — no PII in training."),
    ]
    pw = Inches(3.9)
    px = Inches(0.6)
    py = Inches(1.35)
    for i, (title, desc) in enumerate(pillars):
        row = i // 3
        col_i = i % 3
        x = px + col_i * (pw + Inches(0.28))
        y = py + row * Inches(2.75)
        card(s, x, y, pw, Inches(2.5), border=D100)
        # Icon circle
        R(s, x + Inches(0.18), y + Inches(0.2), Inches(0.48), Inches(0.48), fill=D500)
        TCTR(s, str(i + 1), x + Inches(0.18), y + Inches(0.2), Inches(0.48), Inches(0.48),
             size=14, bold=True, color=WHITE)
        T(s, title, x + Inches(0.8), y + Inches(0.24), pw - Inches(1.0), Inches(0.42),
          size=13, bold=True, color=D700)
        T(s, desc, x + Inches(0.18), y + Inches(0.85), pw - Inches(0.36), Inches(1.45),
          size=11, color=BODY, wrap=True)

    T(s, "SOC 2 Type II and ISO 27001 certification in progress — targeting Q3 2026.",
      Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.35),
      size=11, italic=True, color=MUTED)


# ── SLIDE 10: Pricing ──────────────────────────────────────────────────────────
def s10():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    header_bar(s, "Pricing")
    footer(s)

    plans = [
        ("Free",    "$0",    "month",  D100, D800, False,
         ["50 messages / month", "WhatsApp + Email",
          "Auto-reply AI", "RAG Document Hub", "1 operator"]),
        ("Solo",    "$15",   "month",  D200, D800, False,
         ["200 messages / month", "Auto-reply + RAG",
          "30-day message history", "1 operator"]),
        ("Starter", "$35",   "month",  D500, WHITE, True,
         ["Unlimited messages*", "Smart Escalations",
          "PDPA audit trail", "Up to 3 operators", "Priority support"]),
        ("Pro",     "$55",   "month",  D700, WHITE, False,
         ["Unlimited everything", "Advanced analytics",
          "Custom key phrase maps", "PDPA compliance report",
          "Dedicated support"]),
    ]
    pw = Inches(2.9)
    px = Inches(0.55)
    py = Inches(1.3)
    for name, price, period, col, txt, featured, features in plans:
        ph = Inches(5.6)
        card(s, px, py, pw, ph, fill=col, border=col if featured else D100,
             border_w=Pt(2) if featured else Pt(0.75))
        if featured:
            R(s, px, py, pw, Inches(0.42), fill=D400)
            TCTR(s, "Most Popular", px, py + Inches(0.08), pw, Inches(0.32),
                 size=10, bold=True, color=WHITE)
            ty = py + Inches(0.52)
        else:
            ty = py + Inches(0.22)
        TCTR(s, name, px, ty, pw, Inches(0.55),
             size=22, bold=True, color=txt)
        TCTR(s, price, px, ty + Inches(0.6), pw, Inches(0.75),
             size=40, bold=True, color=txt)
        TCTR(s, f"/{period}", px, ty + Inches(1.32), pw, Inches(0.3),
             size=12, color=txt)
        HR(s, px + Inches(0.28), ty + Inches(1.72), pw - Inches(0.56),
           color=txt if featured else D200, thick=Pt(0.5))
        fy = ty + Inches(1.9)
        for feat in features:
            T(s, f"  {feat}", px + Inches(0.22), fy, pw - Inches(0.38), Inches(0.38),
              size=11, color=txt)
            fy += Inches(0.42)
        px += pw + Inches(0.22)

    T(s, "* Message fair-use limits at 3x plan average. Enterprise pricing on request.",
      Inches(0.55), Inches(6.9), Inches(12.2), Inches(0.3),
      size=10, color=MUTED, italic=True)


# ── SLIDE 11: Getting Started ────────────────────────────────────────────────
def s11():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    header_bar(s, "Getting Started")
    footer(s)

    steps = [
        ("1", "Connect WhatsApp",     "Meta Business Manager → WhatsApp Business API → paste token. ~3 min.", D500),
        ("2", "Connect Email",       "SendGrid Inbound Parse → point webhook at Sequor → add API key. ~5 min.", D600),
        ("3", "Upload Documents",    "Drop PDFs, FAQs, policies into Document Hub. RAG indexes automatically. ~5 min.", D500),
        ("4", "Configure Rules",     "Key phrases → reply strategies → SLA window + escalation contacts. ~10 min.", D600),
        ("5", "Go Live",             "Flip channel to Active. First AI reply fires within the hour.", D500),
    ]
    sx = Inches(0.6)
    sy = Inches(1.35)
    sw = Inches(5.8)
    for num, title, desc, col in steps:
        R(s, sx, sy, Inches(0.55), Inches(0.9), fill=col)
        TCTR(s, num, sx, sy + Inches(0.04), Inches(0.55), Inches(0.9),
             size=18, bold=True, color=WHITE)
        T(s, title, sx + Inches(0.7), sy, sw - Inches(0.7), Inches(0.38),
          size=13, bold=True, color=D800)
        T(s, desc, sx + Inches(0.7), sy + Inches(0.42), sw - Inches(0.7), Inches(0.5),
          size=11, color=BODY, wrap=True)
        sy += Inches(1.05)

    # Right — doc quality note
    card(s, Inches(6.8), Inches(1.35), Inches(5.95), Inches(5.4), border=D100)
    R(s, Inches(6.8), Inches(1.35), Inches(5.95), Inches(0.55), fill=D600)
    T(s, "Document Quality", Inches(7.0), Inches(1.44), Inches(5.6), Inches(0.38),
      size=14, bold=True, color=WHITE)
    doc_tips = [
        ("Best",   "FAQs, policy docs, price lists, service descriptions"),
        ("Good",   "WhatsApp chat exports, organized email threads"),
        ("Difficult", "Verbal knowledge, informal notes, inconsistent docs"),
    ]
    dy = Inches(2.05)
    for label, tip in doc_tips:
        R(s, Inches(7.0), dy + Inches(0.08), Inches(0.08), Inches(0.38),
          fill=D500 if label != "Difficult" else AMBER)
        T(s, label + ":", Inches(7.22), dy, Inches(1.1), Inches(0.38),
          size=11, bold=True, color=D700)
        T(s, tip, Inches(8.35), dy, Inches(4.2), Inches(0.38),
          size=11, color=BODY)
        dy += Inches(0.55)
    T(s, "Start with 5–10 core documents.\nAdd more as you see what the AI answers well.",
      Inches(7.0), dy + Inches(0.2), Inches(5.6), Inches(0.8),
      size=11, italic=True, color=MUTED, wrap=True)


# ── SLIDE 12: Live Demo ────────────────────────────────────────────────────────
def s12():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    header_bar(s, "Live Demo")
    footer(s)

    T(s, "Let's walk through the platform together.",
      Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.5),
      size=16, color=BODY)

    HR(s, Inches(0.6), Inches(1.82), Inches(12.1), color=D100)

    demo_items = [
        ("WhatsApp + Email unified inbox",
         "Both channels in a single thread — no switching between tabs."),
        ("Auto-reply from RAG",
         "Send an email → AI drafts a response from your documents → approve or edit in one tap."),
        ("Smart escalation flow",
         "Flag a thread → assigned with full context → resolution logged. AI learns from it."),
        ("PDPA audit trail",
         "Every action timestamped, tenant-scoped, immutable. Exportable on demand."),
        ("Learning loop",
         "Edit an AI draft → your correction is logged → next time it does better."),
    ]
    dy = Inches(2.05)
    for title, desc in demo_items:
        R(s, Inches(0.6), dy + Inches(0.08), Inches(0.08), Inches(0.42), fill=D500)
        T(s, title, Inches(0.85), dy, Inches(11.6), Inches(0.38),
          size=13, bold=True, color=D700)
        T(s, desc, Inches(0.85), dy + Inches(0.42), Inches(11.6), Inches(0.45),
          size=12, color=BODY, wrap=True)
        dy += Inches(0.98)

    # URL
    R(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5), fill=D050)
    T(s, "🌐  https://badland-swizzle-childlike.ngrok-free.dev",
      Inches(0.85), Inches(6.56), Inches(11.6), Inches(0.38),
      size=13, bold=True, color=D600, align=PP_ALIGN.CENTER)


# ── SLIDE 13: Roadmap ────────────────────────────────────────────────────────
def s13():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    header_bar(s, "Roadmap")
    footer(s)

    quarters = [
        ("Q2 2026", D500, [
            "Portal v1 — WhatsApp + Email",
            "SendGrid + ngrok integration",
            "Learning Loop beta",
            "First lighthouse customers",
        ]),
        ("Q3 2026", D600, [
            "SOC 2 Type II certification",
            "API v1 + webhook outbound",
            "Mobile-responsive portal",
            "Zapier / Make connector",
        ]),
        ("Q4 2026", D700, [
            "Multi-tenant dashboards",
            "Custom key phrase maps",
            "Enterprise tier launch",
            "Thailand + Vietnam expansion",
        ]),
    ]
    rx = Inches(0.6)
    rw = Inches(3.95)
    ry = Inches(1.35)
    for qtr, col, items in quarters:
        card(s, rx, ry, rw, Inches(5.5), border=D100)
        R(s, rx, ry, rw, Inches(0.6), fill=col)
        TCTR(s, qtr, rx, ry + Inches(0.1), rw, Inches(0.5),
             size=18, bold=True, color=WHITE)
        iy = ry + Inches(0.78)
        for item in items:
            R(s, rx + Inches(0.2), iy + Inches(0.1), Inches(0.06), Inches(0.3), fill=col)
            T(s, item, rx + Inches(0.38), iy, rw - Inches(0.56), Inches(0.5),
              size=12, color=BODY, wrap=True)
            iy += Inches(0.6)
        rx += rw + Inches(0.28)


# ── SLIDE 14: Why Sequor ─────────────────────────────────────────────────────
def s14():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    header_bar(s, "Why Sequor")
    footer(s)

    reasons = [
        ("New Category",
         "Not a better version of existing tools. The first platform built for company-wide communication coverage."),
        ("Compounding Moat",
         "Every routing decision, every correction, feeds the learning loop. Gets harder to compete with over time."),
        ("Built for SE Asia",
         "WhatsApp Business is the dominant channel. PDPA compliance from day one. Pricing for SMBs."),
        ("Compliance as Infrastructure",
         "PDPA audit trail, erasure, consent, tenant isolation — not add-ons you pay for. Every tier."),
        ("Setup in Hours",
         "No engineering required. Connect WhatsApp, point email webhook, upload docs. Live same day."),
        ("Institutional Memory",
         "Sequor becomes the system of record for every coverage decision — auditable, searchable."),
    ]
    rw = Inches(3.95)
    rx = Inches(0.6)
    ry = Inches(1.35)
    for title, desc in reasons:
        card(s, rx, ry, rw, Inches(2.5), border=D100)
        R(s, rx, ry, rw, Inches(0.08), fill=D500)
        T(s, title, rx + Inches(0.18), ry + Inches(0.2), rw - Inches(0.36), Inches(0.45),
          size=13, bold=True, color=D700)
        T(s, desc, rx + Inches(0.18), ry + Inches(0.72), rw - Inches(0.36), Inches(1.55),
          size=11, color=BODY, wrap=True)
        ry += Inches(2.62)
        if ry > Inches(5.5):
            ry = Inches(1.35)
            rx += rw + Inches(0.28)


# ── SLIDE 15: Opportunity ─────────────────────────────────────────────────────
def s15():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)
    header_bar(s, "The Opportunity")
    footer(s)

    # Key metrics row
    metrics = [
        ("$2.4B",  "TAM — SE Asia SMB",        "Professional services in Singapore, Thailand,\nVietnam, Malaysia — 5M+ businesses"),
        ("12–18 hrs", "Weekly Coverage Gap",    "Per team of 3–5. Time spent routing, following up,\nand managing what AI should handle."),
        ("$40–60", "Per Seat / Month",         "Professional services teams of 3–10 seats.\n$120–1,000/month ACV. 3–5x productivity ROI."),
        ("8.5/10", "Product Score",             "Problem urgency (9), moat (8), GTM clarity (9),\nregulatory (8), unit economics (8)"),
    ]
    mw = Inches(2.95)
    mx = Inches(0.6)
    my = Inches(1.35)
    for val, lbl, detail in metrics:
        card(s, mx, my, mw, Inches(2.55), border=D100)
        TCTR(s, val, mx, my + Inches(0.18), mw, Inches(0.85),
             size=30, bold=True, color=D500)
        TCTR(s, lbl, mx, my + Inches(1.1), mw, Inches(0.38),
             size=12, bold=True, color=D700)
        T(s, detail, mx + Inches(0.18), my + Inches(1.58), mw - Inches(0.36), Inches(0.82),
          size=10, color=BODY, align=PP_ALIGN.CENTER, wrap=True)
        mx += mw + Inches(0.22)

    # Investment summary
    R(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(2.75), fill=BG_S, line=D100, lw=Pt(0.75))
    T(s, "Investment Summary", Inches(0.85), Inches(4.22), Inches(11.6), Inches(0.5),
      size=18, bold=True, color=D800)
    HR(s, Inches(0.85), Inches(4.78), Inches(11.4), color=D200)

    summary = [
        ("Raise",         "$1.5M Seed"),
        ("Use of Funds",  "Team (40%)  |  Infrastructure (30%)  |  GTM (20%)  |  Compliance (10%)"),
        ("Milestone",     "20 paying customers + $50K MRR in 12 months"),
        ("Runway",        "18 months to Series A gate"),
    ]
    sy = Inches(4.95)
    for label, val in summary:
        T(s, f"{label}:", Inches(0.85), sy, Inches(2.0), Inches(0.42),
          size=13, bold=True, color=D700)
        T(s, val, Inches(2.95), sy, Inches(9.5), Inches(0.42),
          size=13, color=BODY)
        sy += Inches(0.52)


# ── SLIDE 16: Contact ────────────────────────────────────────────────────────
def s16():
    s = prs.slides.add_slide(BLANK)
    bg_white(s)

    # Thin teal bar top
    R(s, 0, 0, W, Inches(0.12), fill=D500)

    # Logo
    R(s, Inches(0.7), Inches(1.5), Inches(1.0), Inches(1.0), fill=D500)
    TCTR(s, "S", Inches(0.7), Inches(1.5), Inches(1.0), Inches(1.0),
         size=42, bold=True, color=WHITE)
    T(s, "Sequor", Inches(2.0), Inches(1.62), Inches(5.0), Inches(0.7),
      size=40, bold=True, color=D800)

    # Horizontal rule
    HR(s, Inches(0.7), Inches(3.0), Inches(5.5), color=D200)

    T(s, "Thank you\nfor your time.",
      Inches(0.7), Inches(3.3), Inches(8.0), Inches(1.4),
      size=34, bold=True, color=D800)
    T(s, "We look forward to building\nthe future of customer communication.",
      Inches(0.7), Inches(4.9), Inches(8.0), Inches(0.85),
      size=16, color=BODY)

    T(s, "May 2026", Inches(0.7), Inches(6.0), Inches(3.0), Inches(0.4),
      size=14, color=MUTED)

    # Right side circles
    R(s, Inches(9.5), Inches(1.2), Inches(3.5), Inches(3.5), fill=D050)
    R(s, Inches(9.8), Inches(1.5), Inches(2.9), Inches(2.9), fill=D200)
    TCTR(s, "AI Coverage\nfor Business",
         Inches(9.8), Inches(1.8), Inches(2.9), Inches(2.2),
         size=20, bold=True, color=D700)

    # Contact info
    R(s, Inches(9.5), Inches(5.1), Inches(3.2), Pt(1), fill=D200)
    contact = ["www.sequor.io", "hello@sequor.io", "Singapore  |  Thailand  |  Vietnam"]
    cy = Inches(5.25)
    for c in contact:
        TCTR(s, c, Inches(9.5), cy, Inches(3.2), Inches(0.4),
             size=12, color=MUTED)
        cy += Inches(0.42)

    footer(s)


# ── Generate ──────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    s1();  s2();  s3();  s4();  s5();  s6()
    s7();  s8();  s9();  s10(); s11(); s12()
    s13(); s14(); s15(); s16()
    prs.save("/Users/aliciapang/Documents/GitHub/Sequor/workspaces/_template/Sequor_Investor_Deck_May2026.pptx")
    print("Saved.")
